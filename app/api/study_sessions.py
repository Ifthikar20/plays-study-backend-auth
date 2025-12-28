"""
Study Sessions API endpoints with AI content processing.
"""
from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from typing import List, Optional, AsyncGenerator
from sqlalchemy.orm import Session
from sqlalchemy import func
from datetime import datetime
from openai import OpenAI
from anthropic import Anthropic
import json
import io
import uuid
import logging
import base64
import tempfile
import subprocess
import os
import hashlib
import asyncio
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

from app.config import settings
from app.dependencies import get_current_active_user, get_db
from app.models.user import User
from app.models.study_session import StudySession
from app.models.topic import Topic
from app.models.question import Question
from app.models.flashcard import Flashcard
from app.models.game_completion import GameCompletion
from app.core.rate_limit import limiter
from app.core.cache import get_cache, set_cache

logger = logging.getLogger(__name__)

router = APIRouter()


def build_topic_hierarchy(session: StudySession, db: Session) -> List["TopicSchema"]:
    """
    Build hierarchical topic structure with questions for a study session.
    OPTIMIZED: Uses pre-loaded relationships when available to avoid N+1 queries.

    Args:
        session: StudySession database model (may have eager-loaded topics)
        db: Database session

    Returns:
        List of TopicSchema objects with nested subtopics and questions
    """
    from app.models.question import Question

    # OPTIMIZED: Check if topics are already loaded (from eager loading)
    # If session.topics is already populated, use it instead of querying
    if hasattr(session, 'topics') and session.topics:
        all_topics = sorted(session.topics, key=lambda t: t.order_index or 0)
    else:
        # Fallback: Query with eager loading for questions
        from sqlalchemy.orm import selectinload
        all_topics = db.query(Topic).options(
            selectinload(Topic.questions)
        ).filter(
            Topic.study_session_id == session.id
        ).order_by(Topic.order_index).all()

    # Build hierarchical structure
    categories = [t for t in all_topics if t.is_category and t.parent_topic_id is None]

    result_topics = []

    for cat_idx, category in enumerate(categories):
        # Get subtopics for this category
        subtopics = [t for t in all_topics if t.parent_topic_id == category.id]

        category_schema = TopicSchema(
            id=f"category-{cat_idx+1}",
            db_id=category.id,
            title=category.title,
            description=category.description or "",
            isCategory=True,
            parentTopicId=None,
            questions=[],
            subtopics=[]
        )

        for sub_idx, subtopic in enumerate(subtopics):
            # OPTIMIZED: Use pre-loaded questions if available
            if hasattr(subtopic, 'questions') and subtopic.questions:
                questions = sorted(subtopic.questions, key=lambda q: q.order_index or 0)
            else:
                # Fallback: Query if not pre-loaded
                questions = db.query(Question).filter(
                    Question.topic_id == subtopic.id
                ).order_by(Question.order_index).all()

            questions_list = [
                QuestionSchema(
                    id=f"topic-{sub_idx+1}-q{q.order_index+1}",
                    question=q.question,
                    options=q.options,
                    correctAnswer=q.correct_answer,
                    explanation=q.explanation,
                    sourceText=q.source_text,
                    sourcePage=q.source_page
                )
                for q in questions
            ]

            subtopic_schema = TopicSchema(
                id=f"subtopic-{cat_idx+1}-{sub_idx+1}",
                db_id=subtopic.id,
                title=subtopic.title,
                description=subtopic.description or "",
                questions=questions_list,
                completed=subtopic.completed or False,
                score=subtopic.score,
                currentQuestionIndex=subtopic.current_question_index or 0,
                isCategory=False,
                parentTopicId=f"category-{cat_idx+1}",
                subtopics=[]
            )
            category_schema.subtopics.append(subtopic_schema)

        result_topics.append(category_schema)

    return result_topics


def analyze_content_complexity(text: str) -> dict:
    """
    Analyze content complexity and recommend topic/question counts.

    Args:
        text: Extracted text content

    Returns:
        Dictionary with analysis results including:
        - word_count: Total number of words
        - estimated_reading_time: Time to read in minutes
        - recommended_topics: Suggested number of topics
        - recommended_questions: Suggested questions per topic
        - complexity_score: 0-1 score of content complexity
        - unique_word_ratio: Vocabulary richness
    """
    words = text.split()
    word_count = len(words)

    # Calculate unique words ratio (vocabulary richness)
    unique_words = len(set(word.lower() for word in words if word.isalnum()))
    unique_word_ratio = unique_words / max(word_count, 1)

    # Average word length (longer words = more complex)
    avg_word_length = sum(len(word) for word in words) / max(word_count, 1)

    # Sentence count (approximate by counting periods, exclamation, question marks)
    sentences = len([c for c in text if c in '.!?'])
    avg_sentence_length = word_count / max(sentences, 1)

    # Calculate complexity score (0-1 scale)
    # Based on: vocabulary richness, average word length, sentence complexity
    complexity_score = min(1.0, (
        (unique_word_ratio * 0.4) +  # 40% weight on vocabulary
        (min(avg_word_length / 8, 1.0) * 0.3) +  # 30% weight on word length
        (min(avg_sentence_length / 25, 1.0) * 0.3)  # 30% weight on sentence length
    ))

    # Reading time (average reading speed: 200-250 words/minute)
    estimated_reading_time = max(1, round(word_count / 225))

    # Recommend topics based on word count and complexity
    # Focus on quality over quantity - only create topics that will have sufficient questions
    # Very short (< 100 words): 1 topic
    # Short (100-500 words): 2-3 topics
    # Medium (500-2000 words): 3-6 topics
    # Long (2000-5000 words): 6-10 topics
    # Very long (5000-10000 words): 10-15 topics
    # Extremely long (10000-20000 words): 15-25 topics
    # Massive (20000+ words): 25-35 topics (maximum to ensure quality)
    if word_count < 100:
        base_topics = 1
    elif word_count < 500:
        base_topics = 2
    elif word_count < 2000:
        base_topics = 4
    elif word_count < 5000:
        base_topics = 8
    elif word_count < 10000:
        base_topics = 12
    elif word_count < 20000:
        base_topics = 20
    else:
        base_topics = 30

    # Adjust based on complexity (more conservative multiplier)
    recommended_topics = max(1, min(35, round(base_topics * (0.8 + complexity_score * 0.4))))

    # Recommend questions per topic based on content depth
    # Aim for maximum questions to ensure comprehensive coverage
    # More questions = better learning coverage and understanding
    if word_count < 1000:
        base_questions = 15
    elif word_count < 3000:
        base_questions = 20
    elif word_count < 10000:
        base_questions = 25
    else:
        base_questions = 30

    # Adjust based on complexity (allow up to 100 questions for complex topics)
    recommended_questions = max(10, min(100, round(base_questions * (0.9 + complexity_score * 0.6))))

    return {
        'word_count': word_count,
        'estimated_reading_time': estimated_reading_time,
        'recommended_topics': recommended_topics,
        'recommended_questions': recommended_questions,
        'complexity_score': round(complexity_score, 2),
        'unique_word_ratio': round(unique_word_ratio, 2),
        'avg_word_length': round(avg_word_length, 1),
        'avg_sentence_length': round(avg_sentence_length, 1)
    }


def detect_file_type_and_extract(content: str) -> tuple[str, str, str]:
    """
    Detect file type and extract text from uploaded content.

    Returns:
        tuple: (extracted_text, file_type, original_base64_content)
            - extracted_text: The text content extracted from the file
            - file_type: One of: 'pdf', 'pptx', 'docx', 'txt'
            - original_base64_content: The original base64 encoded file
    """
    import base64

    content_bytes = None
    file_type = 'txt'
    original_base64 = content

    # First, aggressively try base64 decode (most likely for file uploads)
    try:
        # Try base64 decode - this is the most common format for file uploads via JSON
        decoded = base64.b64decode(content, validate=False)
        # Check if it looks like a valid file (ZIP/docx or PDF)
        if decoded.startswith(b'PK\x03\x04') or decoded.startswith(b'%PDF'):
            content_bytes = decoded
    except Exception:
        pass

    # Process based on file type if we have bytes
    if content_bytes:
        # Office documents (ZIP/docx/pptx)
        if content_bytes.startswith(b'PK\x03\x04'):
            # Try PowerPoint first
            try:
                prs = Presentation(io.BytesIO(content_bytes))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = '\n'.join(text_parts)
                if text.strip():
                    return (text, 'pptx', original_base64)
                raise ValueError("No text content found in PowerPoint")
            except Exception as pptx_error:
                # If PowerPoint fails, try Word document
                try:
                    doc = Document(io.BytesIO(content_bytes))
                    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
                    if text.strip():
                        return (text, 'docx', original_base64)
                    raise ValueError("No text content found in Word document")
                except Exception as docx_error:
                    raise HTTPException(
                        status_code=400,
                        detail=f"Failed to extract text from Office document. PowerPoint error: {str(pptx_error)}. Word error: {str(docx_error)}"
                    )

        # PDF file
        elif content_bytes.startswith(b'%PDF'):
            try:
                pdf_reader = PdfReader(io.BytesIO(content_bytes))
                text = '\n'.join([page.extract_text() for page in pdf_reader.pages])
                if text.strip():
                    return (text, 'pdf', original_base64)
                raise ValueError("No text content found in PDF")
            except Exception as e:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to extract text from PDF: {str(e)}"
                )

    # If we got here, treat as plain text
    cleaned_text = content.replace('\ufffd', '').strip()
    if cleaned_text and len(cleaned_text) > 10:
        return (cleaned_text, 'txt', original_base64)

    # Last resort - return original
    return (content, 'txt', original_base64)


def extract_text_from_content(content: str) -> str:
    """
    Extract readable text from uploaded content.
    Handles:
    - Plain text
    - Word documents (.docx)
    - PowerPoint presentations (.pptx, .ppt)
    - PDF files
    - Base64 encoded content

    Args:
        content: Raw content string (may be binary, text, or base64)

    Returns:
        Extracted text string
    """
    import base64

    content_bytes = None

    # First, aggressively try base64 decode (most likely for file uploads)
    try:
        # Try base64 decode - this is the most common format for file uploads via JSON
        decoded = base64.b64decode(content, validate=False)
        # Check if it looks like a valid file (ZIP/docx or PDF)
        if decoded.startswith(b'PK\x03\x04') or decoded.startswith(b'%PDF'):
            content_bytes = decoded
    except Exception:
        pass

    # If base64 didn't work and content looks like binary, try encoding it
    if content_bytes is None and (content.startswith('PK\x03\x04') or content.startswith('%PDF')):
        # Try different encodings to convert string to bytes
        for encoding in ['latin-1', 'iso-8859-1', 'cp1252', 'utf-8']:
            try:
                content_bytes = content.encode(encoding)
                break
            except (UnicodeEncodeError, UnicodeDecodeError):
                continue

    # Process based on file type if we have bytes
    if content_bytes:
        # Office documents (ZIP/docx/pptx)
        if content_bytes.startswith(b'PK\x03\x04'):
            # Try PowerPoint first
            try:
                prs = Presentation(io.BytesIO(content_bytes))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                text = '\n'.join(text_parts)
                if text.strip():
                    return text
                raise ValueError("No text content found in PowerPoint")
            except Exception as pptx_error:
                # If PowerPoint fails, try Word document
                try:
                    doc = Document(io.BytesIO(content_bytes))
                    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
                    if text.strip():
                        return text
                    raise ValueError("No text content found in Word document")
                except Exception as docx_error:
                    raise HTTPException(
                        status_code=400,
                        detail=f"Failed to extract text from Office document. PowerPoint error: {str(pptx_error)}. Word error: {str(docx_error)}"
                    )

        # PDF file
        elif content_bytes.startswith(b'%PDF'):
            try:
                pdf_reader = PdfReader(io.BytesIO(content_bytes))
                text = '\n'.join([page.extract_text() for page in pdf_reader.pages])
                if text.strip():
                    return text
                raise ValueError("No text content found in PDF")
            except Exception as e:
                raise HTTPException(
                    status_code=400,
                    detail=f"Failed to extract text from PDF: {str(e)}"
                )

    # If we got here, treat as plain text
    # Clean up any Unicode replacement characters
    cleaned_text = content.replace('\ufffd', '').strip()
    if cleaned_text and len(cleaned_text) > 10:
        return cleaned_text

    # Last resort - return original
    return content


def convert_pptx_to_pdf(pptx_base64: str) -> Optional[str]:
    """
    Convert PowerPoint (PPTX) file to PDF using LibreOffice.

    Args:
        pptx_base64: Base64-encoded PPTX file content

    Returns:
        Base64-encoded PDF file content, or None if conversion fails
    """
    temp_dir = None
    try:
        # Create temporary directory for conversion
        temp_dir = tempfile.mkdtemp()
        pptx_path = os.path.join(temp_dir, "presentation.pptx")
        pdf_path = os.path.join(temp_dir, "presentation.pdf")

        # Decode and write PPTX file
        pptx_bytes = base64.b64decode(pptx_base64)
        with open(pptx_path, 'wb') as f:
            f.write(pptx_bytes)

        # Convert using LibreOffice
        # --headless: Run without GUI
        # --convert-to pdf: Convert to PDF format
        # --outdir: Output directory
        try:
            result = subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, pptx_path],
                capture_output=True,
                text=True,
                timeout=30  # 30 second timeout
            )

            # Check if conversion succeeded
            if result.returncode != 0:
                logger.warning(f"LibreOffice conversion failed: {result.stderr}")
                return None

            # Check if PDF file was created
            if not os.path.exists(pdf_path):
                logger.warning("LibreOffice conversion completed but PDF file not found")
                return None

            # Read and encode PDF
            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()

            pdf_base64 = base64.b64encode(pdf_bytes).decode('utf-8')
            logger.info(f"Successfully converted PPTX to PDF ({len(pdf_bytes)} bytes)")
            return pdf_base64

        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice conversion timed out after 30 seconds")
            return None
        except FileNotFoundError:
            logger.warning("LibreOffice not found - PPTX to PDF conversion unavailable")
            return None

    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        return None
    finally:
        # Clean up temporary files
        if temp_dir and os.path.exists(temp_dir):
            try:
                import shutil
                shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"Failed to clean up temp directory: {e}")


class FlashcardSchema(BaseModel):
    """Schema for a flashcard."""
    id: str
    front: str
    back: str
    hint: Optional[str] = None


class QuestionSchema(BaseModel):
    """Schema for a quiz question."""
    id: str
    question: str
    options: List[str]
    correctAnswer: int
    explanation: str
    sourceText: Optional[str] = None  # Source text snippet from document
    sourcePage: Optional[int] = None  # Page number in source document


class TopicSchema(BaseModel):
    """Schema for a study topic."""
    id: str
    db_id: Optional[int] = None  # Database ID for syncing progress
    title: str
    description: str
    questions: List[QuestionSchema] = []
    flashcards: List[FlashcardSchema] = []
    completed: bool = False
    score: Optional[int] = None
    currentQuestionIndex: int = 0
    isCategory: bool = False
    parentTopicId: Optional[str] = None
    subtopics: List['TopicSchema'] = []
    workflowStage: Optional[str] = None

# Enable forward references for recursive schema
TopicSchema.model_rebuild()


class CreateStudySessionRequest(BaseModel):
    """Request schema for creating a study session."""
    title: str = Field(..., min_length=1, max_length=200)
    content: str = Field(..., min_length=10, max_length=100000000)  # 100MB limit for base64 encoded files (large PDFs)
    num_topics: int = Field(default=4, ge=1, le=100)  # Dynamic: 1-100 topics based on content size
    questions_per_topic: int = Field(default=30, ge=5, le=100)  # Generate 30 questions per topic for comprehensive coverage
    progressive_load: bool = Field(default=True)  # ENABLED: Generate questions incrementally using DeepSeek (cheaper) - initial batch only, then call /generate-more-questions


class AnalyzeContentRequest(BaseModel):
    """Request schema for analyzing content before creating a session."""
    content: str = Field(..., min_length=10, max_length=100000000)  # 100MB limit


class ContentAnalysisResponse(BaseModel):
    """Response schema for content analysis."""
    word_count: int
    estimated_reading_time: int  # in minutes
    recommended_topics: int
    recommended_questions: int
    complexity_score: float  # 0-1 scale
    content_summary: str


class CreateStudySessionResponse(BaseModel):
    """Response schema for created study session."""
    id: str  # UUID as string
    title: str
    studyContent: str
    fileContent: Optional[str] = None  # Original file (base64)
    fileType: Optional[str] = None  # File type: pdf, pptx, docx, txt
    pdfContent: Optional[str] = None  # Converted PDF for PPTX files (base64)
    extractedTopics: List[TopicSchema]
    progress: int
    topics: int
    hasFullStudy: bool
    hasSpeedRun: bool
    createdAt: Optional[int] = None  # Unix timestamp in milliseconds
    progressiveLoad: Optional[bool] = False  # Whether questions are loading progressively
    questionsRemaining: Optional[int] = 0  # Number of topics without questions yet
    redirectUrl: str  # Explicit URL for frontend to redirect to after session creation


@router.post("/analyze-content", response_model=ContentAnalysisResponse)
@limiter.limit("10/minute")
async def analyze_content(
    request: Request,
    data: AnalyzeContentRequest,
    current_user: User = Depends(get_current_active_user),
):
    """
    Analyze content and provide recommendations for topics and questions.

    This endpoint analyzes the uploaded content and returns:
    - Word count and estimated reading time
    - Recommended number of topics
    - Recommended questions per topic
    - Complexity score

    Rate Limits:
        - 10 requests per minute per user
    """
    try:
        # Extract text from content (handles Word docs, PDFs, plain text)
        extracted_text = extract_text_from_content(data.content)

        if not extracted_text or len(extracted_text.strip()) < 50:
            raise HTTPException(
                status_code=400,
                detail="Content is too short or empty. Please provide substantial study material (at least 50 characters)."
            )

        # Analyze content complexity
        analysis = analyze_content_complexity(extracted_text)

        # Generate a brief summary (first 200 characters)
        content_preview = extracted_text[:200].strip()
        if len(extracted_text) > 200:
            content_preview += "..."

        return ContentAnalysisResponse(
            word_count=analysis['word_count'],
            estimated_reading_time=analysis['estimated_reading_time'],
            recommended_topics=analysis['recommended_topics'],
            recommended_questions=analysis['recommended_questions'],
            complexity_score=analysis['complexity_score'],
            content_summary=content_preview
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to analyze content: {str(e)}"
        )


@router.post("/create-with-ai", response_model=CreateStudySessionResponse)
@limiter.limit("5/minute")
async def create_study_session_with_ai(
    request: Request,
    data: CreateStudySessionRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Create a new study session and generate topics/questions using AI.

    This endpoint:
    1. Analyzes the provided content
    2. Extracts major topics using AI (dynamically based on content)
    3. Generates quiz questions for each topic
    4. Saves everything to the database

    Rate Limits:
        - 5 requests per minute per user
    """
    try:
        # Validate file size before processing
        # Note: Base64 encoding increases file size by ~33%, so 35MB raw = ~47MB encoded
        MAX_FILE_SIZE_MB = 35
        MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
        content_size = len(data.content)

        if content_size > MAX_FILE_SIZE_BYTES:
            logger.error(f"❌ File too large: {content_size / (1024*1024):.1f}MB (max: {MAX_FILE_SIZE_MB}MB)")
            raise HTTPException(
                status_code=413,
                detail=f"File size ({content_size / (1024*1024):.1f}MB) exceeds maximum allowed size ({MAX_FILE_SIZE_MB}MB). Please use a smaller file or split it into multiple documents."
            )

        logger.info(f"📁 Processing file: {content_size / (1024*1024):.1f}MB")

        # Extract text and detect file type
        extracted_text, file_type, file_content = detect_file_type_and_extract(data.content)

        if not extracted_text or len(extracted_text.strip()) < 50:
            raise HTTPException(
                status_code=400,
                detail="Content is too short or empty. Please provide substantial study material (at least 50 characters)."
            )

        # Convert PPTX to PDF for better rendering in Speed Run mode
        pdf_content = None
        if file_type == 'pptx':
            logger.info("🔄 Converting PowerPoint to PDF for enhanced rendering...")
            pdf_content = convert_pptx_to_pdf(file_content)
            if pdf_content:
                logger.info("✅ PowerPoint successfully converted to PDF")
            else:
                logger.warning("⚠️  PowerPoint to PDF conversion failed - will use extracted text fallback")

        # Check document size and use chunking for large documents
        estimated_tokens = len(extracted_text) // 4  # Rough estimate: 1 token ≈ 4 characters
        logger.info(f"📊 Document size: {len(extracted_text):,} chars, ~{estimated_tokens:,} estimated tokens")

        # Claude 3.5 Haiku has 200k token context window
        # Use chunking for documents that would exceed safe limits
        # Reduced significantly to account for prompt overhead (instructions, subtopics list, output, etc.)
        CHUNK_SIZE_TOKENS = 20000  # 20k tokens per chunk (~80k chars) - conservative to prevent "too large" errors
        OVERLAP_TOKENS = 2000  # 2k token overlap between chunks for context preservation

        chunk_size_chars = CHUNK_SIZE_TOKENS * 4
        overlap_chars = OVERLAP_TOKENS * 4

        # Split document into chunks if necessary
        document_chunks = []
        if estimated_tokens > CHUNK_SIZE_TOKENS:
            logger.info(f"📚 Large document detected. Splitting into chunks of ~{CHUNK_SIZE_TOKENS:,} tokens with {OVERLAP_TOKENS:,} token overlap...")

            # Split into overlapping chunks
            current_pos = 0
            chunk_num = 1
            while current_pos < len(extracted_text):
                end_pos = min(current_pos + chunk_size_chars, len(extracted_text))
                chunk = extracted_text[current_pos:end_pos]
                document_chunks.append(chunk)
                logger.info(f"  📄 Chunk {chunk_num}: {len(chunk):,} chars (~{len(chunk)//4:,} tokens)")

                # Move forward with overlap
                current_pos = end_pos - overlap_chars
                if current_pos >= len(extracted_text) - overlap_chars:
                    break  # Last chunk
                chunk_num += 1

            logger.info(f"✅ Created {len(document_chunks)} chunks for processing")
        else:
            # Document fits in one chunk
            document_chunks = [extracted_text]
            logger.info(f"✅ Document fits in single chunk, no splitting needed")

        # Analyze content to get smart recommendations (use first chunk for analysis)
        analysis = analyze_content_complexity(document_chunks[0])

        # OPTIMIZATION: Cache AI-generated content to save costs on duplicate uploads
        # Generate cache key based on content hash + generation parameters
        content_hash = hashlib.sha256(extracted_text.encode('utf-8')).hexdigest()[:16]
        cache_key = f"ai_session:{content_hash}:{data.num_topics}:{data.questions_per_topic}"

        # Check if we've already generated content for this exact material
        cached_result = get_cache(cache_key)
        if cached_result:
            logger.info(f"✅ Cache HIT - Reusing AI-generated content for hash {content_hash}")
            logger.info(f"💰 Cost savings: Skipped {data.num_topics} topics × {data.questions_per_topic} questions AI generation")

            # Create new study session with cached data
            study_session = StudySession(
                user_id=current_user.id,
                title=cached_result['title'],
                topic=cached_result['topic'],
                study_content=extracted_text,
                file_content=file_content,
                file_type=file_type,
                pdf_content=pdf_content,
                topics_count=cached_result['topics_count'],
                has_full_study=True,
                has_speed_run=True,
                status="in_progress"
            )
            db.add(study_session)
            db.flush()

            # Recreate topics and questions from cached structure
            first_leaf_topic_initialized = False
            for topic_data in cached_result['topics']:
                # Determine initial workflow stage
                # First non-category topic should be quiz_available, others locked
                is_category = topic_data['is_category']
                if not is_category and not first_leaf_topic_initialized:
                    initial_workflow_stage = "quiz_available"
                    first_leaf_topic_initialized = True
                else:
                    initial_workflow_stage = "locked"

                topic = Topic(
                    study_session_id=study_session.id,
                    parent_topic_id=topic_data.get('parent_topic_id'),
                    title=topic_data['title'],
                    description=topic_data['description'],
                    order_index=topic_data['order_index'],
                    is_category=is_category,
                    workflow_stage=initial_workflow_stage
                )
                db.add(topic)
                db.flush()

                # Add questions for this topic
                for q_data in topic_data.get('questions', []):
                    question = Question(
                        topic_id=topic.id,
                        question=q_data['question'],
                        options=q_data['options'],
                        correct_answer=q_data['correct_answer'],
                        explanation=q_data['explanation'],
                        source_text=q_data.get('source_text'),
                        source_page=q_data.get('source_page'),
                        order_index=q_data['order_index']
                    )
                    db.add(question)

                # Add flashcards for this topic (NEW: workflow feature)
                for f_data in topic_data.get('flashcards', []):
                    flashcard = Flashcard(
                        topic_id=topic.id,
                        front=f_data['front'],
                        back=f_data['back'],
                        hint=f_data.get('hint'),
                        order_index=f_data['order_index']
                    )
                    db.add(flashcard)

            db.commit()
            db.refresh(study_session)

            # Build response from cached session
            all_topics = db.query(Topic).filter(
                Topic.study_session_id == study_session.id
            ).order_by(Topic.order_index).all()

            result_topics = build_topic_hierarchy(study_session, db)

            return CreateStudySessionResponse(
                id=str(study_session.id),
                title=study_session.title,
                studyContent=study_session.study_content,
                fileContent=study_session.file_content,
                fileType=study_session.file_type,
                pdfContent=study_session.pdf_content,
                extractedTopics=result_topics,
                progress=0,
                topics=len(result_topics),
                hasFullStudy=True,
                hasSpeedRun=True,
                createdAt=int(study_session.created_at.timestamp() * 1000) if study_session.created_at else None,
                progressiveLoad=False,  # Cached sessions have all questions already
                questionsRemaining=0,  # All questions already generated from cache
                redirectUrl=f"/dashboard/{study_session.id}/full-study"  # Explicit redirect path
            )

        logger.info(f"📝 Cache MISS - Generating new AI content for hash {content_hash}")

        # Progressive loading: For large documents (>5000 words), start with fewer topics
        is_large_doc = analysis['word_count'] > 5000
        initial_topics = data.num_topics

        if data.progressive_load and is_large_doc:
            # Start with only 2-3 categories and 4-6 subtopics for quick initial load
            initial_topics = min(6, data.num_topics)
            logger.info(f"📚 Large document detected ({analysis['word_count']} words). Using progressive load: {initial_topics} initial topics")

        # Calculate number of categories based on topics to generate
        # For 2-5 topics: 2 categories
        # For 6-10 topics: 3 categories
        # For 11-15 topics: 4 categories
        # For 16-20 topics: 5 categories
        num_categories = max(2, min(5, (initial_topics + 3) // 4))
        subtopics_per_category = max(1, initial_topics // num_categories)  # Allow single subtopic per category

        # Initialize AI client (prefer Claude Haiku for speed and quality)
        use_claude = bool(settings.ANTHROPIC_API_KEY)
        anthropic_client = None
        deepseek_client = None

        if use_claude:
            anthropic_client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
            logger.info("🚀 Using Claude Haiku for fast generation")

        # Always initialize DeepSeek as fallback
        if settings.DEEPSEEK_API_KEY:
            deepseek_client = OpenAI(
                api_key=settings.DEEPSEEK_API_KEY,
                base_url="https://api.deepseek.com"
            )
            if not use_claude:
                logger.info("⏱️ Using DeepSeek (consider adding ANTHROPIC_API_KEY for 10x speed)")

        # Step 1: Extract topics from content with hierarchical structure (DYNAMIC PROMPT)
        topics_prompt = f"""Analyze this study material and organize it into a clear, focused hierarchical structure.

Study Material:
{extracted_text}

Content Analysis:
- Word Count: {analysis['word_count']}
- Complexity Score: {analysis['complexity_score']}
- Estimated Reading Time: {analysis['estimated_reading_time']} minutes

Requirements:
1. Create approximately {num_categories} major categories that organize the content at a high level
2. Within each category, break down into focused subtopics (2-3 levels maximum)
3. Prioritize QUALITY over excessive nesting - only create subtopics when there's sufficient content to generate meaningful questions
4. Each leaf node should be a focused concept that can support 25-35 quality questions
5. Provide clear titles and brief descriptions at ALL levels
6. Organize logically (foundational concepts first, building to advanced topics)
7. Create approximately {initial_topics} total LEAF topics across all categories that will actually have questions
8. Avoid creating unnecessary intermediate categories - go directly to testable concepts when possible
9. Focus on core concepts and foundational topics first
10. Remember: The goal is quality questions, not deep nesting

CRITICAL - NO DUPLICATES:
- NEVER create the same topic twice at any level
- Each topic title must be unique within its sibling group
- If a concept appears twice, consolidate it into ONE topic
- Example: Don't create both "Workplace Stress" AND "Work-Related Stress" - they're duplicates!
- Validate each level for duplicate titles before including in JSON

CRITICAL - NO OVERLAPPING:
- Sibling topics must be mutually exclusive (no overlap)
- Each topic should have a distinct, focused scope
- Example GOOD: "Verbal Communication", "Non-Verbal Communication", "Written Communication" (distinct)
- Example BAD: "Face-to-Face Communication", "Verbal Communication", "Body Language" (overlap!)

IMPORTANT: Keep the structure simple and focused. Only nest 2-3 levels deep. Empty subtopics with no questions provide no value.

Return ONLY a valid JSON object in this EXACT format (keep to 2-3 levels maximum):
{{
  "categories": [
    {{
      "title": "Category Title (Level 1)",
      "description": "Brief description of this category",
      "subtopics": [
        {{
          "title": "Subtopic Title (Level 2)",
          "description": "Brief description",
          "subtopics": [
            {{
              "title": "Focused Topic (Level 3 - Leaf)",
              "description": "Brief description of testable concept",
              "subtopics": []
            }}
          ]
        }},
        {{
          "title": "Another Subtopic (Level 2 - Leaf)",
          "description": "Can also be a leaf node if content is focused enough",
          "subtopics": []
        }}
      ]
    }}
  ]
}}

Note: An EMPTY subtopics array [] means this is a LEAF NODE that will have questions generated for it. Keep nesting to 2-3 levels maximum."""

        # Call AI to extract topics (with automatic fallback to DeepSeek if Claude fails)
        topics_text = None
        if use_claude and anthropic_client:
            try:
                topics_response = anthropic_client.messages.create(
                    model="claude-3-5-haiku-20241022",
                    max_tokens=2048,
                    temperature=0.7,
                    messages=[{"role": "user", "content": topics_prompt}]
                )
                topics_text = topics_response.content[0].text
            except Exception as claude_error:
                logger.warning(f"⚠️ Claude API failed: {str(claude_error)}")
                if deepseek_client:
                    logger.info("🔄 Falling back to DeepSeek...")
                    use_claude = False  # Switch to DeepSeek for remaining calls
                else:
                    raise HTTPException(
                        status_code=500,
                        detail=f"Claude API failed and no DeepSeek fallback available: {str(claude_error)}"
                    )

        if not topics_text and deepseek_client:
            topics_response = deepseek_client.chat.completions.create(
                model="deepseek-chat",
                max_tokens=2048,
                temperature=0.7,
                messages=[{"role": "user", "content": topics_prompt}]
            )
            topics_text = topics_response.choices[0].message.content

        if not topics_text:
            raise HTTPException(
                status_code=500,
                detail="No AI provider available. Please configure ANTHROPIC_API_KEY or DEEPSEEK_API_KEY."
            )

        # Parse hierarchical topics
        try:
            start_idx = topics_text.find('{')
            end_idx = topics_text.rfind('}') + 1
            topics_json = json.loads(topics_text[start_idx:end_idx])
            categories_data = topics_json["categories"]
        except (json.JSONDecodeError, KeyError, ValueError) as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse topics: {str(e)}")

        # CRITICAL: Deduplicate topics to prevent duplicate subtopics (e.g., same topic appearing twice)
        def deduplicate_topics(topics_list: list, parent_title: str = "ROOT", level: int = 0) -> tuple:
            """
            Recursively deduplicate topics by title at each level.
            Keeps the first occurrence and removes subsequent duplicates (including all their subtopics).

            Args:
                topics_list: List of topic dictionaries to deduplicate
                parent_title: Title of parent topic for logging context
                level: Current nesting level for logging

            Returns:
                Tuple of (deduplicated_list, duplicate_count)
            """
            seen_titles = set()
            deduplicated = []
            duplicates_removed = 0

            for idx, topic in enumerate(topics_list):
                title = topic.get("title", "").strip()
                title_lower = title.lower()

                # Check if we've already seen this title at this level
                if title_lower in seen_titles:
                    # Count subtopics that will be removed along with duplicate
                    subtopic_count = len(topic.get("subtopics", []))
                    logger.warning(
                        f"{'  ' * level}⚠️ DUPLICATE REMOVED at level {level}: "
                        f"'{title}' (duplicate #{idx + 1}) under '{parent_title}'"
                    )
                    if subtopic_count > 0:
                        logger.warning(
                            f"{'  ' * level}   └─ Also removing {subtopic_count} subtopic(s) under duplicate '{title}'"
                        )
                    duplicates_removed += 1
                    continue

                # Mark this title as seen at this level
                seen_titles.add(title_lower)

                # Recursively deduplicate subtopics
                if "subtopics" in topic and isinstance(topic["subtopics"], list):
                    topic["subtopics"], child_duplicates = deduplicate_topics(
                        topic["subtopics"],
                        parent_title=title,
                        level=level + 1
                    )
                    duplicates_removed += child_duplicates

                deduplicated.append(topic)

            return deduplicated, duplicates_removed

        # Apply deduplication to all categories and their children
        logger.info("🔍 Starting deduplication check across all topic levels...")
        categories_data, total_duplicates = deduplicate_topics(categories_data, parent_title="ROOT", level=0)

        if total_duplicates > 0:
            logger.warning(f"⚠️ REMOVED {total_duplicates} DUPLICATE TOPIC(S) (including their subtopics)")
        else:
            logger.info(f"✅ No duplicates found - all topics are unique")

        logger.info(f"✅ Deduplication complete - {len(categories_data)} unique top-level categories")

        # VERIFICATION: Check for any remaining duplicates across entire hierarchy
        def verify_no_duplicates(topics_list: list, all_titles: dict = None, path: str = "") -> tuple:
            """
            Verify no duplicate titles exist across the entire hierarchy.
            Returns (is_valid, duplicate_info_list)
            """
            if all_titles is None:
                all_titles = {}

            duplicates_found = []

            for topic in topics_list:
                title = topic.get("title", "").strip()
                title_lower = title.lower()
                current_path = f"{path}/{title}" if path else title

                # Check if this title was seen before
                if title_lower in all_titles:
                    duplicates_found.append({
                        "title": title,
                        "first_seen": all_titles[title_lower],
                        "duplicate_at": current_path
                    })
                else:
                    all_titles[title_lower] = current_path

                # Recursively check subtopics
                if "subtopics" in topic and isinstance(topic["subtopics"], list):
                    _, child_duplicates = verify_no_duplicates(
                        topic["subtopics"],
                        all_titles,
                        current_path
                    )
                    duplicates_found.extend(child_duplicates)

            return (len(duplicates_found) == 0, duplicates_found)

        is_valid, found_duplicates = verify_no_duplicates(categories_data)
        if not is_valid:
            logger.error(f"❌ VERIFICATION FAILED: Found {len(found_duplicates)} duplicate(s) after deduplication:")
            for dup in found_duplicates:
                logger.error(f"   - '{dup['title']}' appears at both:")
                logger.error(f"     1. {dup['first_seen']}")
                logger.error(f"     2. {dup['duplicate_at']}")
        else:
            logger.info(f"✅ VERIFICATION PASSED: No duplicates found in entire hierarchy")

        # Count total subtopics for the session
        total_subtopics = sum(len(cat.get("subtopics", [])) for cat in categories_data)

        # Generate a smart title if the provided title is generic (contains "Study Session" and a date)
        session_title = data.title
        if "Study Session" in session_title and any(char.isdigit() for char in session_title):
            # Use the first category as the title for better organization
            if categories_data and len(categories_data) > 0:
                first_category = categories_data[0]["title"]
                # If multiple categories, add a subtitle
                if len(categories_data) > 1:
                    session_title = f"{first_category} + {len(categories_data) - 1} more"
                else:
                    session_title = first_category

        # Step 2: Create study session
        study_session = StudySession(
            user_id=current_user.id,
            title=session_title,
            topic=categories_data[0]["title"] if categories_data else "General Study",
            study_content=extracted_text,  # Store the extracted text
            file_content=file_content,  # Store original file (base64)
            file_type=file_type,  # Store file type (pdf, pptx, docx, txt)
            pdf_content=pdf_content,  # Store converted PDF for PPTX files (base64)
            topics_count=total_subtopics,
            has_full_study=True,
            has_speed_run=True,
            status="in_progress"
        )
        db.add(study_session)
        db.flush()  # Get the session ID

        # Step 3: Create categories and subtopics first, then batch generate ALL questions
        all_topics = []
        subtopic_map = {}  # Map to track subtopics for batch question assignment
        overall_idx = 0
        first_leaf_topic_created = {'value': False}  # Track if first leaf topic has been created (use dict for mutability in closure)

        # Recursive helper function to create topics at any depth
        def create_topics_recursive(
            topic_data: dict,
            parent_topic_id: Optional[int],
            parent_schema: Optional[TopicSchema],
            path: str,  # e.g., "0-1-2" for tracking hierarchy
            order_index: int,
            current_depth: int = 0
        ) -> Optional[TopicSchema]:
            """
            Recursively create topics up to MAX_DEPTH levels.
            Returns the schema for this topic (or None if it's a pure container).
            """
            MAX_DEPTH = 3  # Limit to 3 levels: Category -> Subtopic -> Leaf (or Category -> Leaf)

            subtopics_data = topic_data.get("subtopics", [])

            # Force leaf node if we've reached max depth, regardless of subtopics in data
            if current_depth >= MAX_DEPTH:
                subtopics_data = []
                if topic_data.get("subtopics"):
                    logger.warning(f"⚠️ Max depth {MAX_DEPTH} reached for topic '{topic_data['title']}' - forcing as leaf node")

            has_children = len(subtopics_data) > 0
            is_leaf = not has_children

            # Determine initial workflow stage
            # First leaf topic (non-category) should be quiz_available, all others locked
            if is_leaf and not first_leaf_topic_created['value']:
                initial_workflow_stage = "quiz_available"
                first_leaf_topic_created['value'] = True
                logger.info(f"🎯 First leaf topic '{topic_data['title']}' (path: {path}) - setting to quiz_available")
            else:
                initial_workflow_stage = "locked"

            # SAFETY CHECK: Ensure no duplicate title under same parent (database-level protection)
            existing_sibling = db.query(Topic).filter(
                Topic.study_session_id == study_session.id,
                Topic.parent_topic_id == parent_topic_id,
                Topic.title == topic_data["title"]
            ).first()

            if existing_sibling:
                logger.error(
                    f"❌ DATABASE DUPLICATE DETECTED: Topic '{topic_data['title']}' already exists "
                    f"under parent {parent_topic_id} (session {study_session.id}). Skipping creation."
                )
                return None

            # Create topic in database
            topic = Topic(
                study_session_id=study_session.id,
                parent_topic_id=parent_topic_id,
                title=topic_data["title"],
                description=topic_data.get("description", ""),
                order_index=order_index,
                is_category=has_children,  # Non-leaf nodes are categories
                workflow_stage=initial_workflow_stage
            )
            db.add(topic)
            db.flush()
            logger.debug(f"✅ Created topic '{topic_data['title']}' (ID: {topic.id}, path: {path})")

            # Create schema
            topic_schema = TopicSchema(
                id=f"topic-{path}",
                db_id=topic.id,
                title=topic_data["title"],
                description=topic_data.get("description", ""),
                questions=[],
                completed=False,
                score=None,
                currentQuestionIndex=0,
                isCategory=has_children,
                parentTopicId=parent_schema.id if parent_schema else None,
                subtopics=[]
            )

            # Add ALL topics (both categories and leaf nodes) to subtopic_map for question generation
            # Categories will get overview/synthesis questions, leaf nodes get specific questions
            subtopic_map[path] = {
                "topic": topic,
                "topic_data": topic_data,
                "schema": topic_schema,
                "parent_schema": parent_schema,
                "is_category": has_children,
                "is_leaf": is_leaf
            }

            # Recursively create children
            if has_children:
                for child_idx, child_data in enumerate(subtopics_data):
                    child_path = f"{path}-{child_idx}"
                    child_schema = create_topics_recursive(
                        child_data,
                        topic.id,
                        topic_schema,
                        child_path,
                        child_idx,
                        current_depth + 1  # Increment depth for children
                    )
                    if child_schema:
                        topic_schema.subtopics.append(child_schema)

            return topic_schema

        # First pass: Create all categories and their nested subtopics recursively
        for cat_idx, category_data in enumerate(categories_data):
            category_schema = create_topics_recursive(
                category_data,
                parent_topic_id=None,
                parent_schema=None,
                path=str(cat_idx),
                order_index=cat_idx
            )
            if category_schema:
                all_topics.append(category_schema)

        # FINAL SUMMARY: Report total topics created
        total_topics_in_db = db.query(Topic).filter(Topic.study_session_id == study_session.id).count()
        logger.info(f"📊 FINAL TOPIC COUNT: {total_topics_in_db} total topics created (all unique, no duplicates)")
        logger.info(f"   └─ Top-level categories: {len(all_topics)}")
        logger.info(f"   └─ All subtopics (for questions): {len(subtopic_map)}")

        # Step 4: Generate questions by processing each document chunk
        # For large documents, this processes multiple chunks separately and merges results
        logger.info(f"📡 Processing {len(document_chunks)} document chunk(s) to generate questions...")

        # Get list of all subtopic keys
        all_subtopic_keys = list(subtopic_map.keys())

        # PROGRESSIVE LOADING: Generate questions for only FIRST FEW topics initially
        # This provides fast initial response, then frontend can call /generate-more-questions
        if data.progressive_load:
            # Generate questions for first 2-3 topics only (fast initial load)
            initial_batch_size = min(3, len(all_subtopic_keys))
            subtopics_to_generate = all_subtopic_keys[:initial_batch_size]
            logger.info(f"⚡ PROGRESSIVE LOAD: Generating questions for FIRST {initial_batch_size}/{len(all_subtopic_keys)} subtopics")
            logger.info(f"  Initial batch: {subtopics_to_generate}")
            logger.info(f"  📡 Frontend can call /generate-more-questions to load remaining {len(all_subtopic_keys) - initial_batch_size} topics")
        else:
            # Generate ALL questions upfront (slower but complete)
            subtopics_to_generate = all_subtopic_keys
            logger.info(f"📚 Generating questions for ALL {len(all_subtopic_keys)} subtopics to maximize coverage")
            logger.info(f"  Subtopics: {all_subtopic_keys}")

        # BATCH SUBTOPICS: Process in groups to avoid AI refusing due to response size
        # With 2 subtopics per batch, max ~60 questions per request (2 * 30) - ensures ALL topics get questions
        SUBTOPICS_PER_BATCH = 2  # Reduced to ensure AI generates questions for ALL topics
        subtopic_batches = []

        # Split subtopics into batches (only for selected subtopics_to_generate)
        for i in range(0, len(subtopics_to_generate), SUBTOPICS_PER_BATCH):
            batch_keys = subtopics_to_generate[i:i + SUBTOPICS_PER_BATCH]
            subtopic_batches.append(batch_keys)

        logger.info(f"📦 Split {len(subtopics_to_generate)} subtopics into {len(subtopic_batches)} batches of up to {SUBTOPICS_PER_BATCH}")

        # Collect questions from all chunks
        all_chunk_questions = {}  # {subtopic_key: [questions]}

        # Process each chunk with batched subtopics
        for chunk_idx, chunk_text in enumerate(document_chunks, 1):
            logger.info(f"📄 Processing chunk {chunk_idx}/{len(document_chunks)}...")

            # Process each batch of subtopics for this chunk
            for batch_num, batch_keys in enumerate(subtopic_batches, 1):
                logger.info(f"  📦 Batch {batch_num}/{len(subtopic_batches)}: Processing {len(batch_keys)} subtopics")

                # Build subtopics list for this batch only
                subtopics_list = ""
                for subtopic_key in batch_keys:
                    subtopic_info = subtopic_map[subtopic_key]
                    topic_data = subtopic_info["topic_data"]
                    is_category = subtopic_info.get("is_category", False)
                    is_leaf = subtopic_info.get("is_leaf", True)

                    # Simple format that works for any nesting depth
                    subtopics_list += f"\n[Topic {subtopic_key}]\n"
                    subtopics_list += f"Title: {topic_data['title']}\n"
                    subtopics_list += f"Description: {topic_data.get('description', '')}\n"
                    subtopics_list += f"Type: {'CATEGORY (overview/synthesis questions)' if is_category else 'SPECIFIC TOPIC (detailed questions)'}\n"

                # Build prompt for this chunk and subtopic batch
                batch_prompt = f"""Generate TRICKY and CHALLENGING multiple-choice questions AND flashcards for EACH of the following topics from the study material.

CRITICAL: Generate high-quality, comprehensive questions AND flashcards:
- Extract the MOST IMPORTANT testable concepts, facts, principles, definitions, and examples from the material
- Aim for 25-35 high-quality questions per topic for comprehensive coverage
- Generate 10-15 flashcards per topic for post-quiz review (spaced repetition learning)
- Break down key concepts into multiple questions from different angles
- Test each concept in multiple ways: definition, application, comparison, analysis, synthesis, evaluation
- Create questions for the most important content in the study material
- Generate questions for ALL listed topics below - if the material is light on a topic, use the topic title and description to create questions
- Even if a topic is only mentioned briefly, create comprehensive questions based on what IS mentioned
- For topics not directly covered in the chunk, create questions based on the topic title/description and related content
- Focus on core concepts and ensure complete coverage across all topics

QUESTION TYPES BASED ON TOPIC TYPE:
1. For CATEGORY topics (overview/synthesis questions):
   - Test overall understanding of the category and how subtopics relate
   - Ask synthesis questions that integrate concepts from multiple subtopics
   - Include comparison questions between different subtopics within the category
   - Test the big picture and overarching principles

2. For SPECIFIC TOPIC (detailed questions):
   - Test deep, specific knowledge about that particular topic
   - Include detailed questions about definitions, principles, and mechanisms
   - Ask application questions specific to that topic

3. For ALL topics (CRITICAL - include these):
   - EXAMPLE-BASED QUESTIONS: "Give an example of...", "Which scenario demonstrates...", "Which of the following is an example of..."
   - Application questions that require applying concepts to real scenarios
   - Questions that test practical understanding through examples

DIFFICULTY LEVEL: CHALLENGING
- Make questions that require DEEP analysis and critical thinking
- Use tricky distractors that would fool someone who only skimmed the material
- Test subtle distinctions and nuanced understanding
- Require application of concepts, not just memorization
- Include "all of the above" or "none of the above" when appropriate
- Use comparative questions (e.g., "Which is the PRIMARY..." "What is the MAIN difference...")
- Create questions that test WHY and HOW, not just WHAT
- For every concept, include at least one EXAMPLE-BASED question

Study Material (Chunk {chunk_idx} of {len(document_chunks)}):
{chunk_text}

TOPICS TO COVER ({len(batch_keys)} topics in this batch):
{subtopics_list}

Requirements:
1. Generate 15-20 high-quality questions for EACH topic (total ~30-40 questions for this batch of 2 topics)
2. Generate 8-12 flashcards for EACH topic for spaced repetition review
3. Extract the MOST IMPORTANT testable information from the study material
3. NO DUPLICATES - each question must test a unique concept or angle
4. For EACH topic, include multiple EXAMPLE-BASED questions (e.g., "Which scenario is an example of operant conditioning?")
5. Each question must have exactly 4 PLAUSIBLE options (all should seem correct to someone who doesn't understand deeply)
6. Questions should be TRICKY and CHALLENGING - test deep understanding and critical thinking
7. Distractors should be subtle and based on common misconceptions
8. Provide detailed explanations that explain why the correct answer is right AND why the distractors are wrong
9. For EACH question, include the EXACT source text from the study material with FULL CONTEXT
10. Source text should include the complete sentence(s) or paragraph that contains the answer
11. Include enough surrounding context (2-4 sentences) so students can easily locate it in their document
12. The sourceText must be VERBATIM from the study material - copy it EXACTLY as it appears
13. For PDF documents, estimate the page number where this content appears (if this is chunk {chunk_idx} of {len(document_chunks)}, estimate accordingly)
14. Return ONLY valid JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT

JSON FORMATTING RULES (CRITICAL):
- Use double quotes (") for all strings, not single quotes
- Escape special characters: \" for quotes, \\ for backslashes, \n for newlines
- Do NOT use trailing commas after the last item in arrays or objects
- Ensure all brackets and braces are properly closed
- Numbers for correctAnswer should be integers (0, 1, 2, 3), not strings
- Do NOT include comments in the JSON

IMPORTANT INSTRUCTIONS:
- DO NOT ASK ANY QUESTIONS - you have all the information you need
- DO NOT SAY "I understand" or "I'll help" or provide ANY explanations
- DO NOT ADD ANY TEXT before or after the JSON
- YOUR FIRST CHARACTER MUST BE: {{
- YOUR LAST CHARACTER MUST BE: }}
- OUTPUT ONLY THE JSON OBJECT - NOTHING ELSE

YOU MUST START YOUR RESPONSE WITH THIS EXACT CHARACTER: {{

Begin JSON output now:

GOAL: Create 25-35 comprehensive questions per topic. Focus on QUALITY and coverage of core concepts. Include example-based questions for key concepts.

FLASHCARD REQUIREMENTS:
- Generate 10-15 flashcards per topic for post-quiz review
- Flashcards should cover KEY DEFINITIONS, CONCEPTS, and TERMS
- Front: Question or term (concise, clear)
- Back: Answer or definition (comprehensive but digestible)
- Hint: Optional memory aid or mnemonic (can be null)
- Flashcards complement questions - focus on memorization and quick recall
- Use flashcards for: definitions, formulas, key facts, important dates, terminology

CRITICAL REQUIREMENTS:
- You MUST generate questions AND flashcards for EVERY SINGLE topic key listed above - NO EXCEPTIONS
- If a topic is listed, it MUST appear in your JSON response with questions AND flashcards
- Missing even ONE topic key will result in incomplete learning coverage
- Each topic MUST have 15-20 questions (aim for 20 when possible) AND 8-12 flashcards

Return in this EXACT format (use topic keys EXACTLY as shown above - EVERY topic listed must be in the response):
{{
  "subtopics": {{
    "0": {{
      "questions": [
        {{
          "question": "Question text?",
          "options": ["Option A", "Option B", "Option C", "Option D"],
          "correctAnswer": 0,
          "explanation": "Why this answer is correct",
          "sourceText": "The complete sentence or paragraph from the study material with context.",
          "sourcePage": null
        }},
        ... (25-35 questions for topic "0")
      ],
      "flashcards": [
        {{
          "front": "What is the definition of X?",
          "back": "X is defined as...",
          "hint": "Remember: X starts with..."
        }},
        ... (10-15 flashcards for topic "0")
      ]
    }},
    "0-1": {{
      "questions": [
        {{question object}},
        ... (25-35 questions for topic "0-1")
      ],
      "flashcards": [
        {{flashcard object}},
        ... (10-15 flashcards for topic "0-1")
      ]
    }},
    ... (MUST include ALL topic keys from the list above)
  }}
}}

REMINDER: The response MUST include questions for ALL {len(batch_keys)} topics listed above. Double-check that every topic key appears in your JSON response."""

                # Make API call for this batch
                prompt_tokens = len(batch_prompt) // 4  # Rough estimate
                logger.info(f"    📊 Batch {batch_num} prompt length: {len(batch_prompt):,} characters (~{prompt_tokens:,} tokens)")

                # Note: Removed hard prompt size check - let AI handle it gracefully with error recovery below
                # With reduced chunk size (20k) and batch size (3 subtopics), prompts should be safe

                batch_text = None
                try:
                    if use_claude and anthropic_client:
                        try:
                            # Use prefill technique to force JSON response (prevents conversational responses)
                            batch_response = anthropic_client.messages.create(
                                model="claude-3-5-haiku-20241022",
                                max_tokens=8192,  # Maximum output tokens for Claude 3.5 Haiku
                                temperature=0.7,
                                messages=[
                                    {"role": "user", "content": batch_prompt},
                                    {"role": "assistant", "content": "{"}  # Prefill with opening brace to force JSON
                                ]
                            )
                            # Prepend the opening brace since it was in the prefill
                            batch_text = "{" + batch_response.content[0].text
                        except Exception as claude_error:
                            logger.warning(f"⚠️ Claude API failed for chunk {chunk_idx} batch {batch_num}: {str(claude_error)}")
                            if deepseek_client:
                                logger.info(f"🔄 Falling back to DeepSeek for chunk {chunk_idx} batch {batch_num}...")
                                use_claude = False  # Switch to DeepSeek for remaining calls
                            else:
                                raise  # Re-raise if no fallback available

                    if not batch_text and deepseek_client:
                        batch_response = deepseek_client.chat.completions.create(
                            model="deepseek-chat",
                            max_tokens=8192,  # DeepSeek max limit is 8192
                            temperature=0.7,
                            messages=[{"role": "user", "content": batch_prompt}]
                        )
                        batch_text = batch_response.choices[0].message.content

                    if not batch_text:
                        raise Exception("No AI provider available")

                except Exception as api_error:
                    logger.warning(f"⚠️ AI API call failed for chunk {chunk_idx} batch {batch_num}: {type(api_error).__name__}: {str(api_error)}")
                    # Check if it's a context length error
                    error_msg = str(api_error).lower()
                    if any(keyword in error_msg for keyword in ['context', 'token', 'too long', 'maximum', 'limit']):
                        logger.warning(f"⚠️ Skipping chunk {chunk_idx} batch {batch_num} - too large. Continuing with other batches...")
                        continue  # Skip this batch and continue with next one
                    else:
                        # For non-context errors, still fail (e.g., auth issues, network errors)
                        raise HTTPException(
                            status_code=500,
                            detail=f"AI API error on chunk {chunk_idx} batch {batch_num}: {str(api_error)}"
                        )

                # Log AI response details
                logger.info(f"    📨 Batch {batch_num} - Received AI response, length: {len(batch_text)} characters")

                # Parse batch response with detailed error logging
                try:
                    start_idx = batch_text.find('{')
                    end_idx = batch_text.rfind('}') + 1

                    if start_idx == -1 or end_idx == 0:
                        logger.error(f"❌ Chunk {chunk_idx} Batch {batch_num} - No JSON found in AI response")
                        logger.error(f"❌ AI returned: {batch_text[:500]}...")  # Log first 500 chars
                        chunk_questions = {}
                    else:
                        json_str = batch_text[start_idx:end_idx]
                        batch_json = json.loads(json_str)
                        chunk_questions = batch_json.get("subtopics", {})

                        logger.info(f"    ✅ Batch {batch_num} - Parsed {len(chunk_questions)} subtopics")
                        for key in chunk_questions.keys():
                            q_count = len(chunk_questions[key].get("questions", []))
                            f_count = len(chunk_questions[key].get("flashcards", []))
                            if q_count > 0:
                                logger.info(f"      - Subtopic {key}: {q_count} questions, {f_count} flashcards")

                        # Validate that ALL topics in this batch got questions
                        missing_topics = [key for key in batch_keys if key not in chunk_questions or not chunk_questions[key].get("questions")]
                        if missing_topics:
                            logger.warning(f"    ⚠️ Batch {batch_num} - AI did NOT generate questions for {len(missing_topics)} topics: {missing_topics}")
                            logger.warning(f"    ⚠️ This may result in incomplete coverage. Topics requested: {batch_keys}")
                        else:
                            logger.info(f"    ✅ Batch {batch_num} - ALL {len(batch_keys)} topics have questions")

                except (json.JSONDecodeError, KeyError, ValueError) as e:
                    logger.error(f"❌ Chunk {chunk_idx} Batch {batch_num} - Failed to parse questions: {e}")
                    chunk_questions = {}

                # Merge questions and flashcards from this batch into all_chunk_questions
                for subtopic_key, subtopic_data in chunk_questions.items():
                    questions = subtopic_data.get("questions", [])
                    flashcards = subtopic_data.get("flashcards", [])
                    if questions:
                        if subtopic_key not in all_chunk_questions:
                            all_chunk_questions[subtopic_key] = {"questions": [], "flashcards": []}
                        all_chunk_questions[subtopic_key]["questions"].extend(questions)
                        all_chunk_questions[subtopic_key]["flashcards"].extend(flashcards)
                        logger.debug(f"    🔄 Added {len(questions)} questions and {len(flashcards)} flashcards for subtopic {subtopic_key}")

        # Log merged results
        logger.info(f"🎯 Merging complete - Total subtopics with questions: {len(all_chunk_questions)}")
        for key, data in all_chunk_questions.items():
            q_count = len(data.get("questions", []))
            f_count = len(data.get("flashcards", []))
            logger.info(f"  - Subtopic {key}: {q_count} total questions, {f_count} total flashcards from all chunks")

        # Use merged questions and flashcards as the final subtopics_questions
        subtopics_questions = all_chunk_questions

        # Step 5: Assign questions to subtopics
        question_counter = 0
        logger.info(f"🔄 Assigning questions to {len(subtopic_map)} subtopics...")
        logger.debug(f"🗂️ Subtopic keys in map: {list(subtopic_map.keys())}")
        logger.debug(f"🗂️ Subtopic keys in AI response: {list(subtopics_questions.keys())}")

        for subtopic_key, subtopic_info in subtopic_map.items():
            subtopic = subtopic_info["topic"]
            subtopic_schema = subtopic_info["schema"]
            # Get parent schema (category or parent subtopic)
            parent_schema = subtopic_info.get("parent_schema")

            # Get questions for this subtopic from batch response
            questions_data = subtopics_questions.get(subtopic_key, {}).get("questions", [])

            # Skip if no questions generated (in progressive loading, this is expected for topics beyond the initial batch)
            if not questions_data:
                logger.info(f"⏭️ No questions for subtopic '{subtopic_key}' ('{subtopic.title}') - will be generated later via SSE")
                # Note: The topic is already in the parent's subtopics list (added during creation)
                # It just has an empty questions array, which will be populated later
                continue
            else:
                logger.info(f"✅ Found {len(questions_data)} questions for subtopic '{subtopic_key}' ('{subtopic.title}')")

            # Save ALL questions to database (no limit - variable per topic)
            # Get existing questions for this topic to check for duplicates
            existing_questions = db.query(Question).filter(Question.topic_id == subtopic.id).all()
            existing_question_texts = {q.question.lower().strip() for q in existing_questions}

            questions_list = []
            duplicates_skipped = 0

            for q_idx, q_data in enumerate(questions_data):
                question_text = q_data.get("question", f"Question {q_idx+1}")

                # Check for duplicate questions (case-insensitive)
                if question_text.lower().strip() in existing_question_texts:
                    logger.debug(f"⏭️ Skipping duplicate question: {question_text[:50]}...")
                    duplicates_skipped += 1
                    continue

                # Ensure options is a list
                options = q_data.get("options", ["A", "B", "C", "D"])
                if isinstance(options, str):
                    try:
                        options = json.loads(options)
                    except json.JSONDecodeError:
                        options = ["Option A", "Option B", "Option C", "Option D"]

                question = Question(
                    topic_id=subtopic.id,
                    question=question_text,
                    options=options,
                    correct_answer=q_data.get("correctAnswer", 0),
                    explanation=q_data.get("explanation", ""),
                    source_text=q_data.get("sourceText"),
                    source_page=q_data.get("sourcePage"),
                    order_index=len(questions_list)  # Use actual index in list
                )
                db.add(question)

                # Add to existing set to catch duplicates within this batch
                existing_question_texts.add(question_text.lower().strip())

                questions_list.append(QuestionSchema(
                    id=f"q-{question_counter}",
                    question=q_data.get("question", f"Question {q_idx+1}"),
                    options=options,
                    correctAnswer=q_data.get("correctAnswer", 0),
                    explanation=q_data.get("explanation", ""),
                    sourceText=q_data.get("sourceText"),
                    sourcePage=q_data.get("sourcePage")
                ))
                question_counter += 1

            # Log duplicate detection results
            if duplicates_skipped > 0:
                logger.info(f"  ⏭️ Skipped {duplicates_skipped} duplicate questions for '{subtopic.title}'")

            # Save flashcards to database (NEW: workflow feature)
            flashcards_data = subtopics_questions.get(subtopic_key, {}).get("flashcards", [])
            if flashcards_data:
                logger.info(f"💳 Saving {len(flashcards_data)} flashcards for subtopic '{subtopic.title}'")
                for f_idx, f_data in enumerate(flashcards_data):
                    flashcard = Flashcard(
                        topic_id=subtopic.id,
                        front=f_data.get("front", ""),
                        back=f_data.get("back", ""),
                        hint=f_data.get("hint"),
                        order_index=f_idx
                    )
                    db.add(flashcard)

            # Update subtopic schema with questions
            # Note: subtopic_schema is already in parent's subtopics list (added during creation at line 1212)
            # We're just updating the questions property of the existing schema object
            subtopic_schema.questions = questions_list

        # Validate that questions were generated for a reasonable number of subtopics
        subtopics_with_questions = len([k for k, v in subtopics_questions.items() if v.get("questions")])
        total_subtopics = len(subtopic_map)

        if question_counter == 0:
            logger.error("❌ FATAL: No questions were generated for any subtopic! AI generation completely failed.")
            db.rollback()
            raise HTTPException(
                status_code=500,
                detail="Failed to generate questions. The AI did not return any valid questions. This may be due to document format or content issues. Please try:\n1. A different document\n2. Splitting the document into smaller files\n3. Converting to PDF format if using Word/PowerPoint"
            )

        # Enforce minimum question count for good study sessions
        MIN_QUESTIONS_REQUIRED = 10  # Lowered for nested subtopics that distribute questions across more leaf nodes
        if question_counter < MIN_QUESTIONS_REQUIRED:
            logger.error(f"❌ Only {question_counter} questions generated (minimum: {MIN_QUESTIONS_REQUIRED})")
            db.rollback()
            raise HTTPException(
                status_code=500,
                detail=f"Not enough questions generated ({question_counter}/{MIN_QUESTIONS_REQUIRED} required). The document may be too short, too complex, or in an unsupported format. Please try:\n1. A longer or more detailed document\n2. Converting to PDF format\n3. Checking that the content is educational material"
            )

        # Log coverage statistics
        coverage_percent = (subtopics_with_questions / total_subtopics * 100) if total_subtopics > 0 else 0
        logger.info(f"📊 Question generation coverage: {subtopics_with_questions}/{total_subtopics} subtopics ({coverage_percent:.1f}%)")
        logger.info(f"✅ Successfully generated {question_counter} total questions across {subtopics_with_questions} subtopics")

        # Commit all changes
        db.commit()
        db.refresh(study_session)

        # OPTIMIZATION: Cache the generated content for 24 hours to save on duplicate uploads
        # Store minimal data needed to recreate the session
        try:
            all_topics_from_db = db.query(Topic).filter(
                Topic.study_session_id == study_session.id
            ).order_by(Topic.order_index).all()

            cache_data = {
                'title': study_session.title,
                'topic': study_session.topic,
                'topics_count': study_session.topics_count,
                'topics': []
            }

            # Store topic, question, and flashcard data
            for topic in all_topics_from_db:
                questions = db.query(Question).filter(
                    Question.topic_id == topic.id
                ).order_by(Question.order_index).all()

                flashcards = db.query(Flashcard).filter(
                    Flashcard.topic_id == topic.id
                ).order_by(Flashcard.order_index).all()

                topic_cache = {
                    'title': topic.title,
                    'description': topic.description,
                    'order_index': topic.order_index,
                    'is_category': topic.is_category,
                    'parent_topic_id': topic.parent_topic_id,
                    'questions': [
                        {
                            'question': q.question,
                            'options': q.options,
                            'correct_answer': q.correct_answer,
                            'explanation': q.explanation,
                            'source_text': q.source_text,
                            'source_page': q.source_page,
                            'order_index': q.order_index
                        }
                        for q in questions
                    ],
                    'flashcards': [
                        {
                            'front': f.front,
                            'back': f.back,
                            'hint': f.hint,
                            'order_index': f.order_index
                        }
                        for f in flashcards
                    ]
                }
                cache_data['topics'].append(topic_cache)

            # Cache for 24 hours (86400 seconds)
            set_cache(cache_key, cache_data, ttl=86400)
            logger.info(f"💾 Cached AI-generated content with key {cache_key} (24h TTL)")
            logger.info(f"💰 Future uploads of this content will skip AI generation")
        except Exception as cache_error:
            # Don't fail the request if caching fails
            logger.warning(f"⚠️ Failed to cache generated content: {cache_error}")

        # Calculate how many topics don't have questions yet (for progressive loading)
        topics_without_questions = len(all_subtopic_keys) - len(subtopics_to_generate)

        # Log summary of what's being returned
        logger.info(f"📊 RESPONSE SUMMARY:")
        logger.info(f"  ✅ Generated {question_counter} total questions for {len(subtopics_to_generate)} topics")
        logger.info(f"  📡 {topics_without_questions} topics remaining (will be loaded via SSE)")

        # Count questions in response for verification
        def count_questions_recursive(topic_list):
            total = 0
            for topic in topic_list:
                total += len(topic.questions)
                total += count_questions_recursive(topic.subtopics)
            return total

        questions_in_response = count_questions_recursive(all_topics)
        logger.info(f"  📋 Response includes {questions_in_response} questions in topic tree")

        # Handle both Pydantic model and dict access patterns (defensive coding)
        progressive_load_value = data.progressive_load if hasattr(data, 'progressive_load') else data.get('progressive_load', False)

        return CreateStudySessionResponse(
            id=str(study_session.id),  # Convert UUID to string
            title=study_session.title,
            studyContent=study_session.study_content,
            fileContent=study_session.file_content,
            fileType=study_session.file_type,
            pdfContent=study_session.pdf_content,  # Converted PDF for PPTX files
            extractedTopics=all_topics,
            progress=0,
            topics=len(all_topics),
            hasFullStudy=True,
            hasSpeedRun=True,
            createdAt=int(study_session.created_at.timestamp() * 1000) if study_session.created_at else None,
            progressiveLoad=progressive_load_value,
            questionsRemaining=topics_without_questions,
            redirectUrl=f"/dashboard/{study_session.id}/full-study"  # Explicit redirect path
        )

    except HTTPException:
        # Re-raise HTTP exceptions (like 413 for file size)
        db.rollback()
        raise
    except Exception as api_error:
        logger.error(f"❌ Error in create_study_session_with_ai: {type(api_error).__name__}: {str(api_error)}")
        logger.error(f"❌ Full error: {repr(api_error)}", exc_info=True)
        if "openai" in str(type(api_error).__module__):
            db.rollback()
            raise HTTPException(status_code=500, detail=f"DeepSeek API error: {str(api_error)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"Failed to create study session: {str(api_error)}")


@router.get("/{session_id}", response_model=CreateStudySessionResponse)
async def get_study_session(
    session_id: str,  # UUID as string
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Get a study session with all its topics and questions.

    Returns the complete session data including the hierarchical topic structure
    with all questions, allowing users to resume their study progress.
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'. This may be an old session that is no longer compatible."
        )

    # Fetch session with eager loading of topics and questions
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # Fetch all topics for this session
    all_topics = db.query(Topic).filter(
        Topic.study_session_id == uuid_obj
    ).order_by(Topic.order_index).all()

    # Build hierarchical structure recursively
    def build_topic_tree(topic: Topic, path: str, parent_id: str = None) -> TopicSchema:
        """
        Recursively build topic tree with all questions and subtopics.
        Supports unlimited depth (though we limit to 3 levels during creation).
        """
        # Fetch questions for this topic
        questions = db.query(Question).filter(
            Question.topic_id == topic.id
        ).order_by(Question.order_index).all()

        questions_list = [
            QuestionSchema(
                id=f"{path}-q{q.order_index+1}",
                question=q.question,
                options=q.options,
                correctAnswer=q.correct_answer,
                explanation=q.explanation,
                sourceText=q.source_text,
                sourcePage=q.source_page
            )
            for q in questions
        ]

        # Find all direct children of this topic
        children = [t for t in all_topics if t.parent_topic_id == topic.id]

        # Recursively build subtopics
        subtopics_list = []
        for child_idx, child in enumerate(children):
            child_path = f"{path}-{child_idx+1}"
            child_schema = build_topic_tree(child, child_path, path)
            subtopics_list.append(child_schema)

        # Fetch flashcards for this topic
        flashcards = db.query(Flashcard).filter(
            Flashcard.topic_id == topic.id
        ).order_by(Flashcard.order_index).all()

        flashcards_list = [
            FlashcardSchema(
                id=f"{path}-f{f.order_index+1}",
                front=f.front,
                back=f.back,
                hint=f.hint
            )
            for f in flashcards
        ]

        return TopicSchema(
            id=path,
            db_id=topic.id,
            title=topic.title,
            description=topic.description or "",
            questions=questions_list,
            flashcards=flashcards_list,
            completed=topic.completed or False,
            score=topic.score,
            currentQuestionIndex=topic.current_question_index or 0,
            isCategory=topic.is_category,
            parentTopicId=parent_id,
            subtopics=subtopics_list,
            workflowStage=topic.workflow_stage
        )

    # Start with top-level categories (parent_topic_id=None)
    categories = [t for t in all_topics if t.parent_topic_id is None]

    result_topics = []
    for cat_idx, category in enumerate(categories):
        category_path = f"category-{cat_idx+1}"
        category_schema = build_topic_tree(category, category_path, None)
        result_topics.append(category_schema)

    # Calculate progress
    total_subtopics = sum(len(cat.subtopics) for cat in result_topics)
    completed_subtopics = sum(
        sum(1 for st in cat.subtopics if st.completed)
        for cat in result_topics
    )
    progress = int((completed_subtopics / total_subtopics * 100) if total_subtopics > 0 else 0)

    return CreateStudySessionResponse(
        id=str(session.id),  # Convert UUID to string
        title=session.title,
        studyContent=session.study_content or "",
        fileContent=session.file_content,
        fileType=session.file_type,
        pdfContent=session.pdf_content,  # Converted PDF for PPTX files
        extractedTopics=result_topics,
        progress=progress,
        topics=total_subtopics,
        hasFullStudy=session.has_full_study or False,
        hasSpeedRun=session.has_speed_run or False,
        createdAt=int(session.created_at.timestamp() * 1000) if session.created_at else None
    )


@router.post("/{session_id}/generate-more-questions")
async def generate_more_questions(
    session_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Generate questions for remaining subtopics that don't have questions yet.

    This implements progressive loading - initially only first few subtopics have questions,
    calling this endpoint generates questions for the next batch.

    COST OPTIMIZATION: This endpoint ALWAYS uses DeepSeek API (cheaper) instead of Claude.
    Initial session creation uses Claude (faster/higher quality), but incremental batches
    use DeepSeek to significantly reduce API costs while maintaining quality.
    """
    logger.info(f"📚 Generating more questions for session {session_id}")

    # Validate UUID
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid session ID format")

    # Fetch session
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # OPTIMIZED: Get all topics with question counts in a single aggregated query
    # Subquery to count questions per topic
    question_counts_subquery = db.query(
        Question.topic_id,
        func.count(Question.id).label('question_count')
    ).group_by(Question.topic_id).subquery()

    # Get topics with their question counts
    topics_with_counts = db.query(
        Topic,
        func.coalesce(question_counts_subquery.c.question_count, 0).label('question_count')
    ).outerjoin(
        question_counts_subquery,
        Topic.id == question_counts_subquery.c.topic_id
    ).filter(
        Topic.study_session_id == uuid_obj
    ).order_by(Topic.order_index).all()

    # Find topics without questions using pre-loaded counts
    subtopics_without_questions = [
        topic for topic, count in topics_with_counts if count == 0
    ]

    if not subtopics_without_questions:
        logger.info(f"✅ All subtopics already have questions for session {session_id}")
        return {"message": "All subtopics already have questions", "generated": 0}

    logger.info(f"📊 Found {len(subtopics_without_questions)} topics without questions")

    # Generate questions for next batch (2 topics at a time for better coverage)
    BATCH_SIZE = 2
    next_batch = subtopics_without_questions[:BATCH_SIZE]

    logger.info(f"🔄 Generating questions for next {len(next_batch)} subtopics...")

    # Extract text from stored file content
    if not session.file_content:
        raise HTTPException(status_code=400, detail="No file content available for this session")

    extracted_text, _, _ = detect_file_type_and_extract(session.file_content)

    # Initialize AI client - ALWAYS use DeepSeek for incremental generation (cheaper)
    use_claude = False  # Force DeepSeek for progressive loading to reduce costs
    deepseek_client = OpenAI(
        api_key=settings.DEEPSEEK_API_KEY,
        base_url="https://api.deepseek.com"
    )
    logger.info("⏱️ Using DeepSeek for incremental question generation (cost optimization)")

    # Build prompt for next batch
    subtopics_list = ""
    subtopic_map = {}

    for topic in next_batch:
        # Find parent category
        parent = db.query(Topic).filter(Topic.id == topic.parent_topic_id).first()
        category_title = parent.title if parent else "General"

        # Create mapping key (we'll use topic db_id)
        key = f"topic-{topic.id}"
        subtopic_map[key] = topic

        # Include type information for the AI
        topic_type = "CATEGORY (overview/synthesis questions)" if topic.is_category else "SPECIFIC TOPIC (detailed questions)"

        subtopics_list += f"\n[Topic {key}]\n"
        subtopics_list += f"Title: {topic.title}\n"
        subtopics_list += f"Description: {topic.description or ''}\n"
        subtopics_list += f"Type: {topic_type}\n"
        if parent:
            subtopics_list += f"Parent Category: {category_title}\n"

    # Build prompt
    batch_prompt = f"""Generate TRICKY and CHALLENGING multiple-choice questions AND flashcards for EACH of the following topics from the study material.

CRITICAL: Generate comprehensive questions and flashcards for all topics:
- Extract EVERY testable concept, fact, principle, detail, definition, example, and implication from the material
- Generate 15-20 high-quality questions per topic (reduced from 25-35 to fit within API token limits)
- Generate 8-12 flashcards per topic for spaced repetition learning
- Break down EVERY concept into multiple questions from different angles
- Test each concept in multiple ways: definition, application, comparison, analysis, synthesis, evaluation
- Create questions for every sentence that contains testable information
- Generate questions for ALL listed topics below - if the material is light on a topic, use the topic title and description to create questions
- Even if a topic is only mentioned briefly, create comprehensive questions based on what IS mentioned
- Continue generating until you have exhausted ALL testable content

QUESTION TYPES BASED ON TOPIC TYPE:
1. For CATEGORY topics (overview/synthesis questions):
   - Test overall understanding of the category and how subtopics relate
   - Ask synthesis questions that integrate concepts from multiple subtopics
   - Include comparison questions between different subtopics within the category
   - Test the big picture and overarching principles

2. For SPECIFIC TOPIC (detailed questions):
   - Test deep, specific knowledge about that particular topic
   - Include detailed questions about definitions, principles, and mechanisms
   - Ask application questions specific to that topic

3. For ALL topics (CRITICAL - include these):
   - EXAMPLE-BASED QUESTIONS: "Give an example of...", "Which scenario demonstrates...", "Which of the following is an example of..."
   - Application questions that require applying concepts to real scenarios
   - Questions that test practical understanding through examples

DIFFICULTY LEVEL: CHALLENGING
- Make questions that require DEEP analysis and critical thinking
- Use tricky distractors that would fool someone who only skimmed the material
- Test subtle distinctions and nuanced understanding
- Require application of concepts, not just memorization
- Include "all of the above" or "none of the above" when appropriate
- Use comparative questions (e.g., "Which is the PRIMARY..." "What is the MAIN difference...")
- Create questions that test WHY and HOW, not just WHAT
- For every concept, include at least one EXAMPLE-BASED question

Study Material:
{extracted_text[:80000]}

TOPICS TO COVER ({len(next_batch)} topics in this batch):
{subtopics_list}

Requirements:
1. Generate 15-20 high-quality questions for EACH topic (total ~30-40 questions for this batch of 2 topics)
2. Generate 8-12 flashcards for EACH topic for spaced repetition review
3. Extract the MOST IMPORTANT testable information from the study material
3. NO DUPLICATES - each question must test a unique concept or angle
4. For EACH topic, include multiple EXAMPLE-BASED questions (e.g., "Which scenario is an example of...?")
5. Each question must have exactly 4 PLAUSIBLE options (all should seem correct to someone who doesn't understand deeply)
6. Questions should be TRICKY and CHALLENGING - test deep understanding and critical thinking
7. Distractors should be subtle and based on common misconceptions
8. Provide detailed explanations that explain why the correct answer is right AND why the distractors are wrong
9. For EACH question, include the EXACT source text from the study material with FULL CONTEXT
10. Source text should include the complete sentence(s) or paragraph that contains the answer
11. Include enough surrounding context (2-4 sentences) so students can easily locate it in their document
12. The sourceText must be VERBATIM from the study material - copy it EXACTLY as it appears
13. For PDF documents, estimate which section/page the content appears in
14. Return ONLY valid JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT

JSON FORMATTING RULES (CRITICAL):
- Use double quotes (") for all strings, not single quotes
- Escape special characters: \" for quotes, \\ for backslashes, \n for newlines
- Do NOT use trailing commas after the last item in arrays or objects
- Ensure all brackets and braces are properly closed
- Numbers for correctAnswer should be integers (0, 1, 2, 3), not strings
- Do NOT include comments in the JSON

IMPORTANT INSTRUCTIONS:
- DO NOT ASK ANY QUESTIONS - you have all the information you need
- DO NOT SAY "I understand" or "I'll help" or provide ANY explanations
- DO NOT ADD ANY TEXT before or after the JSON
- YOUR FIRST CHARACTER MUST BE: {{
- YOUR LAST CHARACTER MUST BE: }}
- OUTPUT ONLY THE JSON OBJECT - NOTHING ELSE

YOU MUST START YOUR RESPONSE WITH THIS EXACT CHARACTER: {{

Begin JSON output now:

GOAL: Create 15-20 comprehensive questions AND 8-12 flashcards per topic. Focus on QUALITY and coverage of core concepts. Include example-based questions for key concepts.

FLASHCARD REQUIREMENTS:
- Generate 8-12 flashcards per topic for post-quiz review
- Flashcards should cover KEY DEFINITIONS, CONCEPTS, and TERMS
- Front: Question or term (concise, clear)
- Back: Answer or definition (comprehensive but digestible)
- Hint: Optional memory aid or mnemonic (can be null)
- Flashcards complement questions - focus on memorization and quick recall

CRITICAL REQUIREMENTS:
- You MUST generate questions AND flashcards for EVERY SINGLE topic key listed above - NO EXCEPTIONS
- If a topic is listed, it MUST appear in your JSON response with questions AND flashcards
- Missing even ONE topic key will result in incomplete learning coverage
- Each topic MUST have 15-20 questions (aim for 20 when possible) AND 8-12 flashcards

Return in this EXACT format (use topic keys EXACTLY as shown above - EVERY topic listed must be in the response):
{{
  "subtopics": {{
    "topic-123": {{
      "questions": [
        {{
          "question": "Question text?",
          "options": ["Option A", "Option B", "Option C", "Option D"],
          "correctAnswer": 0,
          "explanation": "Why this answer is correct",
          "sourceText": "The complete sentence or paragraph from the study material with context.",
          "sourcePage": null
        }},
        ... (15-20 questions for topic "topic-123")
      ],
      "flashcards": [
        {{
          "front": "What is the definition of X?",
          "back": "X is defined as...",
          "hint": "Remember: X starts with..."
        }},
        ... (8-12 flashcards for topic "topic-123")
      ]
    }},
    "topic-456": {{
      "questions": [
        {{question object}},
        ... (15-20 questions for topic "topic-456")
      ],
      "flashcards": [
        {{flashcard object}},
        ... (8-12 flashcards for topic "topic-456")
      ]
    }},
    ... (MUST include ALL topic keys from the list above)
  }}
}}

REMINDER: The response MUST include questions AND flashcards for ALL {len(next_batch)} topics listed above. Double-check that every topic key appears in your JSON response."""

    # Make API call
    logger.info(f"📊 Prompt length: {len(batch_prompt):,} characters")

    try:
        if use_claude:
            # Use prefill technique to force JSON response (prevents conversational responses)
            batch_response = anthropic_client.messages.create(
                model="claude-3-5-haiku-20241022",
                max_tokens=8192,  # Maximum for Haiku (2 topics × ~20 questions each fits in 8k)
                temperature=0.7,
                messages=[
                    {"role": "user", "content": batch_prompt},
                    {"role": "assistant", "content": "{"}  # Prefill with opening brace to force JSON
                ]
            )
            # Prepend the opening brace since it was in the prefill
            batch_text = "{" + batch_response.content[0].text
            logger.info(f"📊 AI response stats: stop_reason={batch_response.stop_reason}, input_tokens={batch_response.usage.input_tokens}, output_tokens={batch_response.usage.output_tokens}")
        else:
            batch_response = deepseek_client.chat.completions.create(
                model="deepseek-chat",
                max_tokens=8192,  # DeepSeek max limit is 8192
                temperature=0.7,
                messages=[{"role": "user", "content": batch_prompt}]
            )
            batch_text = batch_response.choices[0].message.content

            # Check for truncation in DeepSeek response
            finish_reason = batch_response.choices[0].finish_reason
            logger.info(f"📊 DeepSeek response stats: finish_reason={finish_reason}, prompt_tokens={batch_response.usage.prompt_tokens}, completion_tokens={batch_response.usage.completion_tokens}")

            if finish_reason == "length":
                logger.error(f"❌ DeepSeek response TRUNCATED - hit max_tokens limit!")
                logger.error(f"❌ Response was cut off at {batch_response.usage.completion_tokens} tokens")
                logger.error(f"❌ Last 500 chars of truncated response: ...{batch_text[-500:]}")
                raise HTTPException(
                    status_code=500,
                    detail=f"AI response was truncated at token limit. Please reduce the number of topics or questions per batch."
                )
    except HTTPException:
        # Re-raise HTTP exceptions as-is
        raise
    except Exception as api_error:
        logger.error(f"❌ AI API call failed: {type(api_error).__name__}: {str(api_error)}")
        raise HTTPException(status_code=500, detail=f"AI API error: {str(api_error)}")

    # Parse response
    logger.info(f"📨 Received AI response, length: {len(batch_text)} characters")

    try:
        start_idx = batch_text.find('{')
        end_idx = batch_text.rfind('}') + 1

        if start_idx == -1 or end_idx == 0:
            logger.error(f"❌ No JSON found in AI response")
            logger.error(f"❌ Full AI response (first 2000 chars):")
            logger.error(f"{batch_text[:2000]}")
            raise HTTPException(status_code=500, detail=f"AI returned non-JSON response. First 500 chars: {batch_text[:500]}")

        json_str = batch_text[start_idx:end_idx]

        # Log the JSON string for debugging
        logger.debug(f"📝 Attempting to parse JSON (length: {len(json_str)} chars)")
        logger.debug(f"📝 First 500 chars of JSON: {json_str[:500]}")

        batch_json = json.loads(json_str)
        subtopics_questions = batch_json.get("subtopics", {})

        logger.info(f"✅ Parsed {len(subtopics_questions)} topics from AI response")
        for key in subtopics_questions.keys():
            q_count = len(subtopics_questions[key].get("questions", []))
            logger.info(f"  - Topic {key}: {q_count} questions")

        # Validate that ALL topics in this batch got questions
        missing_topics = [key for key in subtopic_map.keys() if key not in subtopics_questions or not subtopics_questions[key].get("questions")]
        if missing_topics:
            logger.warning(f"⚠️ AI did NOT generate questions for {len(missing_topics)} topics: {missing_topics}")
            logger.warning(f"⚠️ Topics requested: {list(subtopic_map.keys())}")
        else:
            logger.info(f"✅ ALL {len(subtopic_map)} topics have questions")

    except json.JSONDecodeError as e:
        logger.error(f"❌ JSON parsing failed: {e}")
        logger.error(f"❌ Error at line {e.lineno}, column {e.colno}")
        logger.error(f"❌ Full JSON response (first 3000 chars):")
        logger.error(f"{json_str[:3000]}")
        # Show context around the error
        if e.pos and e.pos < len(json_str):
            start = max(0, e.pos - 200)
            end = min(len(json_str), e.pos + 200)
            logger.error(f"❌ Context around error position {e.pos}:")
            logger.error(f"{json_str[start:end]}")
        raise HTTPException(status_code=500, detail=f"Failed to parse AI JSON: {str(e)}")
    except (KeyError, ValueError) as e:
        logger.error(f"❌ Failed to parse questions: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse AI response: {str(e)}")

    # Save questions and flashcards to database
    total_questions_generated = 0
    total_flashcards_generated = 0

    for key, topic in subtopic_map.items():
        questions_data = subtopics_questions.get(key, {}).get("questions", [])
        flashcards_data = subtopics_questions.get(key, {}).get("flashcards", [])

        if not questions_data:
            logger.warning(f"⚠️ No questions generated for subtopic '{key}' ('{topic.title}')")
            continue

        logger.info(f"✅ Saving {len(questions_data)} questions for subtopic '{topic.title}'")

        # Save questions
        for q_idx, q_data in enumerate(questions_data):
            # Ensure options is a list
            options = q_data.get("options", ["A", "B", "C", "D"])
            if isinstance(options, str):
                try:
                    options = json.loads(options)
                except json.JSONDecodeError:
                    options = ["Option A", "Option B", "Option C", "Option D"]

            question = Question(
                topic_id=topic.id,
                question=q_data.get("question", f"Question {q_idx+1}"),
                options=options,
                correct_answer=q_data.get("correctAnswer", 0),
                explanation=q_data.get("explanation", ""),
                source_text=q_data.get("sourceText"),
                source_page=q_data.get("sourcePage"),
                order_index=q_idx
            )
            db.add(question)
            total_questions_generated += 1

        # Save flashcards
        if flashcards_data:
            logger.info(f"💳 Saving {len(flashcards_data)} flashcards for subtopic '{topic.title}'")
            for f_idx, f_data in enumerate(flashcards_data):
                flashcard = Flashcard(
                    topic_id=topic.id,
                    front=f_data.get("front", ""),
                    back=f_data.get("back", ""),
                    hint=f_data.get("hint"),
                    order_index=f_idx
                )
                db.add(flashcard)
                total_flashcards_generated += 1

    # Commit to database
    try:
        db.commit()
        logger.info(f"✅ Successfully saved {total_questions_generated} questions and {total_flashcards_generated} flashcards to database")
    except Exception as e:
        db.rollback()
        logger.error(f"❌ Failed to save questions: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save questions: {str(e)}")

    # Count remaining subtopics without questions
    remaining_count = len(subtopics_without_questions) - len(next_batch)

    logger.info(f"✅ Generated {total_questions_generated} questions and {total_flashcards_generated} flashcards for {len(next_batch)} subtopics")
    logger.info(f"📊 Remaining subtopics without questions: {remaining_count}")

    return {
        "message": f"Successfully generated questions for {len(next_batch)} subtopics",
        "generated": len(next_batch),
        "totalQuestions": total_questions_generated,
        "totalFlashcards": total_flashcards_generated,
        "remaining": remaining_count,
        "hasMore": remaining_count > 0
    }


@router.get("/{session_id}/generate-more-questions-stream")
async def generate_more_questions_stream(
    session_id: str,
    token: str,  # Token from query param (EventSource doesn't support headers)
    db: Session = Depends(get_db),
):
    """
    Stream question generation progress using Server-Sent Events (SSE).

    This endpoint streams real-time progress updates as questions are generated
    for remaining topics without questions. Uses DeepSeek API exclusively for
    cost optimization.

    SSE Events:
    - start: Initial connection with total topics remaining
    - batch_start: Before each batch generation
    - progress: After each batch completes (with detailed stats)
    - complete: When all questions are generated
    - error: If any error occurs
    """

    async def event_generator() -> AsyncGenerator[str, None]:
        """Generate Server-Sent Events as questions are created"""

        try:
            # Authenticate user from token (EventSource doesn't support headers)
            from app.core.security import decode_access_token

            try:
                payload = decode_access_token(token)
                user_id = int(payload.get("sub"))

                current_user = db.query(User).filter(User.id == user_id).first()
                if not current_user or not current_user.is_active:
                    yield f"event: error\ndata: {json.dumps({'error': 'Authentication failed'})}\n\n"
                    return
            except Exception as auth_error:
                logger.error(f"❌ SSE authentication failed: {auth_error}")
                yield f"event: error\ndata: {json.dumps({'error': 'Invalid authentication token'})}\n\n"
                return

            # Validate UUID
            try:
                uuid_obj = uuid.UUID(session_id)
            except ValueError:
                yield f"event: error\ndata: {json.dumps({'error': 'Invalid session ID format'})}\n\n"
                return

            # Fetch session
            session = db.query(StudySession).filter(
                StudySession.id == uuid_obj,
                StudySession.user_id == current_user.id
            ).first()

            if not session:
                yield f"event: error\ndata: {json.dumps({'error': 'Study session not found'})}\n\n"
                return

            # Find topics without questions (optimized query)
            question_counts_subquery = db.query(
                Question.topic_id,
                func.count(Question.id).label('question_count')
            ).group_by(Question.topic_id).subquery()

            topics_with_counts = db.query(
                Topic,
                func.coalesce(question_counts_subquery.c.question_count, 0).label('question_count')
            ).outerjoin(
                question_counts_subquery,
                Topic.id == question_counts_subquery.c.topic_id
            ).filter(
                Topic.study_session_id == uuid_obj
            ).order_by(Topic.order_index).all()

            subtopics_without_questions = [
                topic for topic, count in topics_with_counts if count == 0
            ]

            if not subtopics_without_questions:
                logger.info(f"✅ All subtopics already have questions for session {session_id}")
                yield f"event: complete\ndata: {json.dumps({'message': 'All topics already have questions', 'generated': 0, 'remaining': 0})}\n\n"
                return

            total_remaining = len(subtopics_without_questions)
            logger.info(f"📊 Found {total_remaining} topics without questions")

            # Send start event
            yield f"event: start\ndata: {json.dumps({'totalRemaining': total_remaining, 'sessionId': session_id})}\n\n"

            # Extract text from file (needed for AI generation)
            if not session.file_content:
                yield f"event: error\ndata: {json.dumps({'error': 'No file content available'})}\n\n"
                return

            extracted_text, _, _ = detect_file_type_and_extract(session.file_content)

            # Initialize DeepSeek client (ALWAYS use DeepSeek for cost optimization)
            deepseek_client = OpenAI(
                api_key=settings.DEEPSEEK_API_KEY,
                base_url="https://api.deepseek.com"
            )
            logger.info("⏱️ Using DeepSeek for incremental question generation (SSE stream)")

            # Generate in batches
            BATCH_SIZE = 2
            batch_num = 0
            total_questions_all_batches = 0
            total_flashcards_all_batches = 0

            while subtopics_without_questions:
                next_batch = subtopics_without_questions[:BATCH_SIZE]
                batch_num += 1

                logger.info(f"🔄 SSE Stream - Batch {batch_num}: Generating for {len(next_batch)} topics")

                # Send batch start event
                batch_start_data = {'batchNumber': batch_num, 'topicsInBatch': len(next_batch)}
                yield f"event: batch_start\ndata: {json.dumps(batch_start_data)}\n\n"

                # Build prompt for this batch
                subtopics_list = ""
                subtopic_map = {}

                for topic in next_batch:
                    parent = db.query(Topic).filter(Topic.id == topic.parent_topic_id).first()
                    category_title = parent.title if parent else "General"
                    key = f"topic-{topic.id}"
                    subtopic_map[key] = topic
                    topic_type = "CATEGORY (overview/synthesis questions)" if topic.is_category else "SPECIFIC TOPIC (detailed questions)"

                    subtopics_list += f"\n[Topic {key}]\n"
                    subtopics_list += f"Title: {topic.title}\n"
                    subtopics_list += f"Description: {topic.description or ''}\n"
                    subtopics_list += f"Type: {topic_type}\n"
                    if parent:
                        subtopics_list += f"Parent Category: {category_title}\n"

                # Build AI prompt (same as generate-more-questions)
                batch_prompt = f"""Generate TRICKY and CHALLENGING multiple-choice questions AND flashcards for EACH of the following topics from the study material.

CRITICAL: Generate comprehensive questions and flashcards for all topics:
- Extract EVERY testable concept, fact, principle, detail, definition, example, and implication from the material
- Generate 15-20 high-quality questions per topic (reduced from 25-35 to fit within API token limits)
- Generate 8-12 flashcards per topic for spaced repetition learning

Study Material:
{extracted_text[:80000]}

TOPICS TO COVER ({len(next_batch)} topics in this batch):
{subtopics_list}

Requirements:
1. Generate 15-20 high-quality questions for EACH topic (total ~30-40 questions for this batch of 2 topics)
2. Generate 8-12 flashcards for EACH topic for spaced repetition review
3. Each question must have exactly 4 PLAUSIBLE options
4. Provide detailed explanations
5. Include source text from the study material
6. Return ONLY valid JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT

IMPORTANT INSTRUCTIONS:
- DO NOT ASK ANY QUESTIONS - you have all the information you need
- DO NOT SAY "I understand" or "I'll help" or provide ANY explanations
- YOUR FIRST CHARACTER MUST BE: {{
- OUTPUT ONLY THE JSON OBJECT - NOTHING ELSE

Return in this EXACT format:
{{
  "subtopics": {{
    "topic-123": {{
      "questions": [/* 25-35 questions */],
      "flashcards": [/* 10-15 flashcards */]
    }}
  }}
}}"""

                try:
                    # Call DeepSeek API
                    batch_response = deepseek_client.chat.completions.create(
                        model="deepseek-chat",
                        max_tokens=8192,
                        temperature=0.7,
                        messages=[{"role": "user", "content": batch_prompt}]
                    )
                    batch_text = batch_response.choices[0].message.content

                    # Parse JSON response
                    start_idx = batch_text.find('{')
                    end_idx = batch_text.rfind('}') + 1

                    if start_idx == -1 or end_idx == 0:
                        logger.error(f"❌ No JSON found in AI response")
                        error_data = {'error': 'AI returned non-JSON response', 'batch': batch_num}
                        yield f"event: error\ndata: {json.dumps(error_data)}\n\n"
                        break

                    json_str = batch_text[start_idx:end_idx]
                    batch_json = json.loads(json_str)
                    subtopics_questions = batch_json.get("subtopics", {})

                    # Save questions and flashcards
                    total_questions_generated = 0
                    total_flashcards_generated = 0

                    for key, topic in subtopic_map.items():
                        questions_data = subtopics_questions.get(key, {}).get("questions", [])
                        flashcards_data = subtopics_questions.get(key, {}).get("flashcards", [])

                        # Save questions
                        for q_idx, q_data in enumerate(questions_data):
                            options = q_data.get("options", ["A", "B", "C", "D"])
                            if isinstance(options, str):
                                try:
                                    options = json.loads(options)
                                except:
                                    options = ["Option A", "Option B", "Option C", "Option D"]

                            question = Question(
                                topic_id=topic.id,
                                question=q_data.get("question", f"Question {q_idx+1}"),
                                options=options,
                                correct_answer=q_data.get("correctAnswer", 0),
                                explanation=q_data.get("explanation", ""),
                                source_text=q_data.get("sourceText"),
                                source_page=q_data.get("sourcePage"),
                                order_index=q_idx
                            )
                            db.add(question)
                            total_questions_generated += 1

                        # Save flashcards
                        if flashcards_data:
                            for f_idx, f_data in enumerate(flashcards_data):
                                flashcard = Flashcard(
                                    topic_id=topic.id,
                                    front=f_data.get("front", ""),
                                    back=f_data.get("back", ""),
                                    hint=f_data.get("hint"),
                                    order_index=f_idx
                                )
                                db.add(flashcard)
                                total_flashcards_generated += 1

                    db.commit()
                    logger.info(f"💾 Database committed - {total_questions_generated}Q and {total_flashcards_generated}F saved")

                    # Log which topics got questions
                    for key, topic in subtopic_map.items():
                        q_count = len(subtopics_questions.get(key, {}).get("questions", []))
                        f_count = len(subtopics_questions.get(key, {}).get("flashcards", []))
                        logger.info(f"   📝 Topic ID {topic.id} '{topic.title}': {q_count}Q, {f_count}F")

                    # Update counters
                    total_questions_all_batches += total_questions_generated
                    total_flashcards_all_batches += total_flashcards_generated

                    # Remove processed topics
                    subtopics_without_questions = subtopics_without_questions[BATCH_SIZE:]
                    remaining_count = len(subtopics_without_questions)

                    # Send progress event
                    progress_data = {
                        'batchNumber': batch_num,
                        'generated': len(next_batch),
                        'remaining': remaining_count,
                        'totalQuestions': total_questions_generated,
                        'totalFlashcards': total_flashcards_generated,
                        'cumulativeQuestions': total_questions_all_batches,
                        'cumulativeFlashcards': total_flashcards_all_batches,
                        'hasMore': remaining_count > 0
                    }

                    logger.info(f"✅ SSE Stream - Batch {batch_num} complete: {total_questions_generated}Q, {total_flashcards_generated}F. Remaining: {remaining_count}")
                    logger.info(f"📡 Sending progress event to frontend: {progress_data}")
                    yield f"event: progress\ndata: {json.dumps(progress_data)}\n\n"

                    # Small delay between batches
                    await asyncio.sleep(0.5)

                except Exception as e:
                    logger.error(f"❌ SSE Stream - Batch {batch_num} failed: {e}")
                    error_data = {'error': str(e), 'batch': batch_num}
                    yield f"event: error\ndata: {json.dumps(error_data)}\n\n"
                    break

            # Send completion event
            completion_data = {
                'message': 'All questions generated successfully',
                'totalQuestions': total_questions_all_batches,
                'totalFlashcards': total_flashcards_all_batches,
                'batchesCompleted': batch_num
            }
            logger.info(f"🎉 SSE Stream complete: {total_questions_all_batches}Q, {total_flashcards_all_batches}F")
            yield f"event: complete\ndata: {json.dumps(completion_data)}\n\n"

        except Exception as e:
            import traceback
            logger.error(f"❌ SSE stream error: {e}")
            logger.error(f"❌ Traceback: {traceback.format_exc()}")
            error_data = {'error': str(e)}
            yield f"event: error\ndata: {json.dumps(error_data)}\n\n"

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",  # Disable nginx buffering
            "Connection": "keep-alive",
        }
    )


@router.delete("/{session_id}")
async def delete_study_session(
    session_id: str,  # UUID as string
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Delete a study session and all its associated data.

    This will cascade delete all topics and questions associated with the session.
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'."
        )

    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    db.delete(session)
    db.commit()

    return {"message": "Study session deleted successfully", "id": session_id}


@router.patch("/{session_id}/archive")
async def archive_study_session(
    session_id: str,  # UUID as string
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Archive a study session.

    Changes the session status to 'archived'.
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'."
        )

    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    session.status = "archived"
    db.commit()
    db.refresh(session)

    return {"message": "Study session archived successfully", "id": session_id}


class UpdateTopicProgressRequest(BaseModel):
    """Request schema for updating topic progress."""
    score: int = Field(..., ge=0, le=100)  # Score as percentage 0-100
    current_question_index: int = Field(..., ge=0)
    completed: bool = False


class UpdateUserXPRequest(BaseModel):
    """Request schema for updating user XP."""
    xp_to_add: int = Field(..., ge=0, le=1000)  # XP to add (capped at 1000 per call)


@router.patch("/{session_id}/topics/{topic_id}/progress")
@limiter.limit("60/minute")
async def update_topic_progress(
    request: Request,
    session_id: str,
    topic_id: int,  # Database topic ID
    data: UpdateTopicProgressRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Update topic progress (score, current question index, completion status).

    This endpoint is called after each answer to persist user progress.

    Rate Limits:
        - 60 requests per minute per user
    """
    # Validate UUID format
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        logger.error(f"❌ Invalid session ID format: {session_id} (expected UUID)")
        raise HTTPException(
            status_code=400,
            detail=f"Invalid session ID format. Expected UUID, got: '{session_id}'."
        )

    # Verify session belongs to user
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # Find the topic
    topic = db.query(Topic).filter(
        Topic.id == topic_id,
        Topic.study_session_id == uuid_obj
    ).first()

    if not topic:
        raise HTTPException(status_code=404, detail="Topic not found")

    # Update topic progress
    topic.score = data.score
    topic.current_question_index = data.current_question_index
    topic.completed = data.completed

    # Handle workflow stage transitions
    # When quiz is completed, transition to flashcard_review stage
    if data.completed and topic.workflow_stage == "quiz_available":
        topic.workflow_stage = "flashcard_review"
        logger.info(f"✅ Quiz completed for topic {topic_id} ({topic.title}) - transitioning to flashcard_review stage")

    # Update session progress (percentage of completed subtopics)
    all_topics = db.query(Topic).filter(
        Topic.study_session_id == uuid_obj,
        Topic.is_category == False  # Only count leaf topics
    ).all()

    completed_topics = sum(1 for t in all_topics if t.completed)
    total_topics = len(all_topics)
    session.progress = int((completed_topics / total_topics * 100) if total_topics > 0 else 0)

    db.commit()
    db.refresh(topic)

    return {
        "message": "Topic progress updated successfully",
        "topic_id": topic_id,
        "score": topic.score,
        "current_question_index": topic.current_question_index,
        "completed": topic.completed,
        "workflow_stage": topic.workflow_stage,
        "session_progress": session.progress
    }


@router.patch("/user/xp")
@limiter.limit("100/minute")
async def update_user_xp(
    request: Request,
    data: UpdateUserXPRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Add XP to the current user's total.

    This endpoint is called after each correct answer.

    Rate Limits:
        - 100 requests per minute per user
    """
    # Update user XP
    current_user.xp += data.xp_to_add
    current_user.updated_at = datetime.utcnow()

    # Calculate level (simple formula: level = floor(xp / 100) + 1)
    # Every 100 XP = 1 level
    current_user.level = (current_user.xp // 100) + 1

    db.commit()
    db.refresh(current_user)

    return {
        "message": "XP updated successfully",
        "xp": current_user.xp,
        "xp_added": data.xp_to_add,
        "level": current_user.level
    }


class BatchXPUpdate(BaseModel):
    """OPTIMIZATION: Batch multiple XP updates to reduce API calls by 90%."""
    xp_increments: List[int] = Field(..., min_items=1, max_items=100)  # Batch up to 100 answers


@router.post("/user/xp/batch")
@limiter.limit("20/minute")  # Much lower limit since it's batched
async def batch_update_user_xp(
    request: Request,
    data: BatchXPUpdate,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    OPTIMIZATION: Add XP in batch (every 5-10 answers instead of every 1).

    This reduces API calls by 90% and database writes by 90%.

    Example:
        Instead of calling /user/xp 10 times with xp_to_add=10,
        call /user/xp/batch once with xp_increments=[10,10,10,10,10,10,10,10,10,10]

    Rate Limits:
        - 20 requests per minute per user (much lower because it's batched)
    """
    total_xp = sum(data.xp_increments)

    # Update user XP
    current_user.xp += total_xp
    current_user.updated_at = datetime.utcnow()

    # Calculate level (simple formula: level = floor(xp / 100) + 1)
    current_user.level = (current_user.xp // 100) + 1

    db.commit()
    db.refresh(current_user)

    return {
        "message": f"Batch updated {len(data.xp_increments)} XP increments",
        "answers_processed": len(data.xp_increments),
        "total_xp_added": total_xp,
        "xp": current_user.xp,
        "level": current_user.level
    }


class TopicProgressUpdate(BaseModel):
    """Single topic progress update."""
    topic_id: int
    score: int = Field(..., ge=0, le=100)
    current_question_index: int = Field(..., ge=0)
    completed: bool = False


class BatchProgressUpdate(BaseModel):
    """OPTIMIZATION: Batch multiple topic progress updates to reduce API calls by 90%."""
    session_id: str  # UUID
    updates: List[TopicProgressUpdate] = Field(..., min_items=1, max_items=50)


@router.post("/batch-progress")
@limiter.limit("30/minute")  # Lower limit since it's batched
async def batch_update_progress(
    request: Request,
    data: BatchProgressUpdate,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    OPTIMIZATION: Update multiple topic progress in one request.

    This reduces API calls by 90% when students complete multiple questions.

    Example:
        Instead of calling /topics/{id}/progress 10 times,
        call /batch-progress once with all 10 updates

    Rate Limits:
        - 30 requests per minute per user (much lower because it's batched)
    """
    # Validate UUID
    try:
        uuid_obj = uuid.UUID(data.session_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid session ID format")

    # Verify session belongs to user
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # Process all updates
    updated_topics = []
    for update in data.updates:
        # Find and update topic
        topic = db.query(Topic).filter(
            Topic.id == update.topic_id,
            Topic.study_session_id == uuid_obj
        ).first()

        if topic:
            topic.score = update.score
            topic.current_question_index = update.current_question_index
            topic.completed = update.completed
            updated_topics.append(topic.id)
        else:
            logger.warning(f"⚠️ Topic {update.topic_id} not found in session {data.session_id}")

    # Update session progress (percentage of completed subtopics)
    all_topics = db.query(Topic).filter(
        Topic.study_session_id == uuid_obj,
        Topic.is_category == False  # Only count leaf topics
    ).all()

    completed_topics = sum(1 for t in all_topics if t.completed)
    total_topics = len(all_topics)
    session.progress = int((completed_topics / total_topics * 100) if total_topics > 0 else 0)

    db.commit()

    return {
        "message": f"Batch updated {len(updated_topics)} topics",
        "topics_updated": len(updated_topics),
        "topics_completed": completed_topics,
        "total_topics": total_topics,
        "session_progress": session.progress,
        "updated_topic_ids": updated_topics
    }


# ===== FLASHCARD ENDPOINTS (Workflow Feature) =====

@router.get("/topics/{topic_id}/flashcards")
async def get_topic_flashcards(
    topic_id: int,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Get all flashcards for a topic (for post-quiz review).

    This endpoint returns flashcards that should be reviewed AFTER
    completing the quiz for a topic. Flashcards use spaced repetition
    for long-term memorization.

    Returns:
        List of flashcards with spaced repetition metadata
    """
    # Verify topic exists and belongs to user's session
    topic = db.query(Topic).join(StudySession).filter(
        Topic.id == topic_id,
        StudySession.user_id == current_user.id
    ).first()

    if not topic:
        raise HTTPException(status_code=404, detail="Topic not found")

    # Get all flashcards for this topic
    flashcards = db.query(Flashcard).filter(
        Flashcard.topic_id == topic_id
    ).order_by(Flashcard.order_index).all()

    return {
        "topic_id": topic_id,
        "topic_title": topic.title,
        "flashcards": [
            {
                "id": f.id,
                "front": f.front,
                "back": f.back,
                "hint": f.hint,
                "order_index": f.order_index,
                "ease_factor": f.ease_factor,
                "interval_days": f.interval_days,
                "repetitions": f.repetitions,
                "next_review_date": f.next_review_date.isoformat() if f.next_review_date else None,
                "last_reviewed_at": f.last_reviewed_at.isoformat() if f.last_reviewed_at else None,
                "total_reviews": f.total_reviews,
                "correct_reviews": f.correct_reviews,
                "accuracy": f.get_accuracy(),
                "is_due": f.is_due_for_review()
            }
            for f in flashcards
        ],
        "total_flashcards": len(flashcards),
        "due_for_review": sum(1 for f in flashcards if f.is_due_for_review())
    }


class FlashcardReviewRequest(BaseModel):
    """Request schema for submitting a flashcard review."""
    quality: int = Field(..., ge=0, le=5, description="Rating 0-5 (0=failed, 3=hard, 4=good, 5=easy)")


@router.post("/flashcards/{flashcard_id}/review")
async def submit_flashcard_review(
    flashcard_id: int,
    data: FlashcardReviewRequest,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Submit a review for a flashcard (updates spaced repetition algorithm).

    Quality ratings:
    - 0-2: Failed (card will be shown again soon)
    - 3: Hard (passed but difficult)
    - 4: Good (passed with some effort)
    - 5: Easy (passed easily)

    The algorithm automatically calculates the next review date based
    on your performance (SM-2 spaced repetition algorithm).

    Returns:
        Updated flashcard with new review schedule
    """
    # Verify flashcard exists and belongs to user's session
    flashcard = db.query(Flashcard).join(Topic).join(StudySession).filter(
        Flashcard.id == flashcard_id,
        StudySession.user_id == current_user.id
    ).first()

    if not flashcard:
        raise HTTPException(status_code=404, detail="Flashcard not found")

    # Update spaced repetition data
    flashcard.calculate_next_review(data.quality)

    db.commit()
    db.refresh(flashcard)

    return {
        "flashcard_id": flashcard.id,
        "quality": data.quality,
        "next_review_date": flashcard.next_review_date.isoformat() if flashcard.next_review_date else None,
        "interval_days": flashcard.interval_days,
        "ease_factor": flashcard.ease_factor,
        "repetitions": flashcard.repetitions,
        "total_reviews": flashcard.total_reviews,
        "correct_reviews": flashcard.correct_reviews,
        "accuracy": flashcard.get_accuracy()
    }


@router.post("/topics/{topic_id}/flashcards/complete")
async def mark_flashcards_complete(
    topic_id: int,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Mark all flashcards for a topic as completed.

    This endpoint is called when the user finishes reviewing all flashcards
    for a topic. It:
    1. Updates the topic's workflow_stage to "completed"
    2. Unlocks next topics (topics that have this topic as a prerequisite)

    This is the final step in the workflow: quiz → flashcards → complete

    Returns:
        Success status and list of newly unlocked topics
    """
    # Verify topic exists and belongs to user's session
    topic = db.query(Topic).join(StudySession).filter(
        Topic.id == topic_id,
        StudySession.user_id == current_user.id
    ).first()

    if not topic:
        raise HTTPException(status_code=404, detail="Topic not found")

    # Update topic workflow stage to completed
    topic.workflow_stage = "completed"
    topic.completed = True

    logger.info(f"✅ Topic {topic_id} ({topic.title}) marked as completed (flashcards done)")

    # Unlock next topics (topics that have this topic as a prerequisite)
    unlocked_topics = []

    # Find all topics that have this topic in their prerequisites
    dependent_topics = db.query(Topic).filter(
        Topic.study_session_id == topic.study_session_id,
        Topic.prerequisite_topic_ids.contains([topic_id])  # PostgreSQL array contains
    ).all()

    for dependent in dependent_topics:
        # Check if ALL prerequisites are now completed
        if dependent.prerequisite_topic_ids:
            prerequisites = db.query(Topic).filter(
                Topic.id.in_(dependent.prerequisite_topic_ids)
            ).all()

            all_prerequisites_completed = all(
                prereq.workflow_stage == "completed"
                for prereq in prerequisites
            )

            if all_prerequisites_completed and dependent.workflow_stage == "locked":
                # Unlock this topic!
                dependent.workflow_stage = "quiz_available"
                unlocked_topics.append({
                    "topic_id": dependent.id,
                    "title": dependent.title
                })
                logger.info(f"🔓 Unlocked topic {dependent.id} ({dependent.title}) - all prerequisites completed")

    db.commit()

    return {
        "success": True,
        "topic_id": topic_id,
        "title": topic.title,
        "workflow_stage": "completed",
        "unlocked_topics": unlocked_topics,
        "unlocked_count": len(unlocked_topics)
    }


# ===== WORKFLOW VISUALIZATION ENDPOINT =====

@router.get("/sessions/{session_id}/workflow")
async def get_session_workflow(
    session_id: str,
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db),
):
    """
    Get workflow visualization data for a study session.

    Returns topic hierarchy with:
    - Position coordinates (for skill-tree UI)
    - Workflow stages (locked, quiz_available, quiz_completed, flashcard_review, completed)
    - Prerequisites and dependencies
    - Question and flashcard counts
    - Progress tracking

    This data is used by the frontend to render a visual skill-tree
    showing the user's learning path and progress.
    """
    # Validate UUID
    try:
        uuid_obj = uuid.UUID(session_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid session ID format")

    # Verify session belongs to user
    session = db.query(StudySession).filter(
        StudySession.id == uuid_obj,
        StudySession.user_id == current_user.id
    ).first()

    if not session:
        raise HTTPException(status_code=404, detail="Study session not found")

    # Get all topics with question and flashcard counts
    topics = db.query(
        Topic,
        func.count(Question.id.distinct()).label('question_count'),
        func.count(Flashcard.id.distinct()).label('flashcard_count')
    ).outerjoin(Question).outerjoin(Flashcard).filter(
        Topic.study_session_id == uuid_obj
    ).group_by(Topic.id).all()

    # Get all game completions for this user and session's topics
    topic_ids = [topic.id for topic, _, _ in topics]
    game_completions = db.query(GameCompletion).filter(
        GameCompletion.user_id == current_user.id,
        GameCompletion.topic_id.in_(topic_ids)
    ).all()

    # Create a map of topic_id -> list of game completions
    games_by_topic = {}
    for gc in game_completions:
        if gc.topic_id not in games_by_topic:
            games_by_topic[gc.topic_id] = []
        games_by_topic[gc.topic_id].append({
            "game_type": gc.game_type,
            "completed": gc.completed,
            "score": gc.score
        })

    # Build workflow nodes (topics + games)
    workflow_nodes = []
    game_nodes = []

    for topic, question_count, flashcard_count in topics:
        workflow_stage = topic.workflow_stage or "locked"

        # Helper fields for frontend (derived from workflow_stage)
        quiz_completed = workflow_stage in ["quiz_completed", "flashcard_review", "completed"]
        flashcards_completed = workflow_stage == "completed"

        # Games are recommended when flashcards are completed
        games_available = flashcards_completed

        node = {
            "node_type": "topic",  # Distinguish from game nodes
            "topic_id": topic.id,
            "title": topic.title,
            "description": topic.description,
            "is_category": topic.is_category,
            "parent_topic_id": topic.parent_topic_id,
            "order_index": topic.order_index,
            "workflow_stage": workflow_stage,
            "position_x": topic.position_x,
            "position_y": topic.position_y,
            "prerequisite_topic_ids": topic.prerequisite_topic_ids or [],
            "question_count": question_count,
            "flashcard_count": flashcard_count,
            "completed": topic.completed,
            "score": topic.score,
            "current_question_index": topic.current_question_index,

            # Helper fields for frontend React Flow implementation
            "quiz_completed": quiz_completed,  # True if quiz is done (flashcards now available)
            "flashcards_completed": flashcards_completed,  # True if flashcards are done (topic fully completed)
            "games_available": games_available,  # True if games are recommended

            # Game completion info for this topic
            "games": games_by_topic.get(topic.id, [])
        }
        workflow_nodes.append(node)

        # Add game nodes for completed topics
        if games_available and not topic.is_category:
            # Recommend primary games (Memory Match and True/False)
            primary_games = [
                {
                    "game_type": "memory_match",
                    "title": "🧠 Memory Match",
                    "description": "Match concepts and definitions",
                    "icon": "🧠"
                },
                {
                    "game_type": "true_false",
                    "title": "✓ True or False",
                    "description": "Test your knowledge with scenarios",
                    "icon": "✓"
                }
            ]

            for game_info in primary_games:
                # Check if this game has been completed
                completed_game = next(
                    (g for g in games_by_topic.get(topic.id, []) if g['game_type'] == game_info['game_type']),
                    None
                )

                game_node = {
                    "node_type": "game",
                    "game_type": game_info['game_type'],
                    "title": game_info['title'],
                    "description": game_info['description'],
                    "icon": game_info['icon'],
                    "topic_id": topic.id,  # Link back to parent topic
                    "topic_title": topic.title,
                    "completed": completed_game['completed'] if completed_game else False,
                    "score": completed_game['score'] if completed_game and completed_game['completed'] else None,
                    "workflow_stage": "available",  # Games are always available once unlocked
                }
                game_nodes.append(game_node)

    return {
        "session_id": str(session.id),
        "title": session.title,
        "progress": session.progress,
        "workflow_nodes": workflow_nodes,  # Topic nodes
        "game_nodes": game_nodes,  # Game nodes (separate for easier frontend rendering)
        "total_nodes": len(workflow_nodes),
        "total_games": len(game_nodes)
    }
