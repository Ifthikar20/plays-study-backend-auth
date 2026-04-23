"""
AI service — content extraction, analysis, and topic/question generation.
Ported from FastAPI study_sessions.py to a clean Django service.
"""
import io
import json
import logging
import base64
from django.conf import settings

logger = logging.getLogger(__name__)


# ─── Text Extraction ──────────────────────────────────────

def extract_text(content: str) -> str:
    """Extract readable text from content (plain text, base64 PDF/DOCX/PPTX)."""
    _, _, text = detect_file_type(content)
    return text


def detect_file_type(content: str) -> tuple:
    """
    Detect file type and extract text.

    The frontend always sends the uploaded file as a base64 string
    (stripped of the `data:` prefix), including for plain-text files.
    We try binary formats first via magic bytes, then fall back to
    decoding the base64 as UTF-8 text. If the content isn't base64 at
    all (e.g. pasted content from the API), we return it as-is.

    Returns: (extracted_text, file_type, original_content)
    """
    content = (content or '').strip()
    if not content:
        return ('', 'txt', content)

    decoded: bytes | None = None
    try:
        # validate=False lets us accept base64 with whitespace, but we then
        # re-validate by round-tripping: only treat as base64 if the decoded
        # bytes aren't just a noisy reinterpretation of the input.
        candidate = base64.b64decode(content, validate=False)
        # Reject trivial decodes (e.g. short plain text that happens to decode)
        if len(candidate) > 4:
            decoded = candidate
    except Exception:
        decoded = None

    if decoded:
        # PDF
        if decoded.startswith(b'%PDF'):
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(decoded))
                text = '\n'.join(page.extract_text() or '' for page in reader.pages)
                if text.strip():
                    return (text, 'pdf', content)
            except Exception as exc:
                logger.warning('PDF parse failed: %s', exc)

        # Office (Zip container) — try PowerPoint then Word
        if decoded.startswith(b'PK\x03\x04'):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(decoded))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            parts.append(shape.text)
                text = '\n'.join(parts)
                if text.strip():
                    return (text, 'pptx', content)
            except Exception:
                pass
            try:
                from docx import Document
                doc = Document(io.BytesIO(decoded))
                text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
                if text.strip():
                    return (text, 'docx', content)
            except Exception as exc:
                logger.warning('DOCX parse failed: %s', exc)

        # Plain UTF-8 text hiding inside a base64 blob
        try:
            text = decoded.decode('utf-8')
            if text.strip() and _looks_like_text(text):
                return (text, 'txt', content)
        except UnicodeDecodeError:
            pass

    # Not base64 — treat the raw string as text directly.
    cleaned = content.replace('\ufffd', '').strip()
    return (cleaned if cleaned else content, 'txt', content)


def _looks_like_text(s: str) -> bool:
    """Heuristic: is this decoded blob readable plain text rather than binary?"""
    if not s:
        return False
    # Reject if a large fraction of characters are outside common printable range.
    printable = sum(1 for c in s if c.isprintable() or c in '\n\r\t')
    return (printable / len(s)) > 0.9


# ─── Complexity Analysis ──────────────────────────────────

def analyze_complexity(text: str) -> dict:
    """Analyze content complexity and recommend topic/question counts."""
    words = text.split()
    word_count = len(words)
    unique_words = len(set(w.lower() for w in words if w.isalnum()))
    unique_ratio = unique_words / max(word_count, 1)
    avg_word_len = sum(len(w) for w in words) / max(word_count, 1)
    sentences = len([c for c in text if c in '.!?'])
    avg_sent_len = word_count / max(sentences, 1)

    complexity = min(1.0, unique_ratio * 0.4 + min(avg_word_len / 8, 1.0) * 0.3 + min(avg_sent_len / 25, 1.0) * 0.3)

    if word_count < 500:
        topics = 2
    elif word_count < 2000:
        topics = 4
    elif word_count < 5000:
        topics = 8
    else:
        topics = min(20, word_count // 500)

    questions = max(10, min(30, int(15 * (0.9 + complexity * 0.6))))

    return {
        'word_count': word_count,
        'estimated_reading_time': max(1, round(word_count / 225)),
        'recommended_topics': topics,
        'recommended_questions': questions,
        'complexity_score': round(complexity, 2),
    }


# ─── AI Topic/Question Generation ─────────────────────────

def generate_topics_and_questions(text: str, num_topics: int, questions_per_topic: int) -> list:
    """
    Generate topics, flashcards, and questions from study content.

    Tries Anthropic first, then OpenAI, and finally the deterministic placeholder
    generator. Any provider-level failure (network, model retired, malformed JSON)
    falls through silently — the user always gets a usable session, even if every
    AI provider is misconfigured.
    """
    if settings.ANTHROPIC_API_KEY:
        try:
            topics = _generate_with_anthropic(text, num_topics, questions_per_topic)
            if topics:
                return topics
            logger.warning('Anthropic returned no topics — falling back')
        except Exception as exc:
            logger.warning('Anthropic generation failed (%s) — falling back', exc)

    if settings.OPENAI_API_KEY:
        try:
            topics = _generate_with_openai(text, num_topics, questions_per_topic)
            if topics:
                return topics
            logger.warning('OpenAI returned no topics — falling back')
        except Exception as exc:
            logger.warning('OpenAI generation failed (%s) — falling back', exc)

    logger.info('Using deterministic placeholder topic generator')
    return _generate_placeholder(text, num_topics, questions_per_topic)


def _generate_with_anthropic(text: str, num_topics: int, qpt: int) -> list:
    """Generate using Claude Haiku 4.5 — fast and cheap for structured JSON."""
    from anthropic import Anthropic

    client = Anthropic(api_key=settings.ANTHROPIC_API_KEY)
    prompt = _build_prompt(text, num_topics, qpt)

    # Haiku 4.5 — current model (replaces the retired claude-3-5-haiku-latest).
    # Cheap enough at $1/$5 per 1M tokens to use freely; max 64K output tokens.
    response = client.messages.create(
        model='claude-haiku-4-5',
        max_tokens=8000,
        messages=[{'role': 'user', 'content': prompt}],
    )

    return _parse_ai_response(response.content[0].text)


def _generate_with_openai(text: str, num_topics: int, qpt: int) -> list:
    """Generate using OpenAI as a backup provider."""
    from openai import OpenAI

    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    prompt = _build_prompt(text, num_topics, qpt)

    response = client.chat.completions.create(
        model='gpt-4o-mini',
        messages=[{'role': 'user', 'content': prompt}],
        max_tokens=8000,
    )

    return _parse_ai_response(response.choices[0].message.content)


def _build_prompt(text: str, num_topics: int, qpt: int) -> str:
    """Build the AI prompt for topic/question/flashcard generation."""
    # Truncate text if too long
    max_chars = 60000
    if len(text) > max_chars:
        text = text[:max_chars] + '\n\n[Content truncated for processing]'

    fpt = max(4, min(8, qpt // 2))  # flashcards per subtopic

    return f"""Analyze the following study material and create a structured learning plan.

Create {num_topics} main topic categories, each with 2-3 subtopics.
For each subtopic, generate:
  - {qpt} multiple-choice questions (with exactly 4 options each, one correct)
  - {fpt} flashcards, each with a concise prompt (front) and a clear, self-contained answer (back).
    Flashcards MUST be derivable from the source material. The "front" is the prompt, the "back" is the answer.

Return ONLY valid JSON in this exact format (no markdown, no prose before or after):
[
  {{
    "title": "Category Title",
    "description": "Brief description",
    "subtopics": [
      {{
        "title": "Subtopic Title",
        "description": "Brief description",
        "flashcards": [
          {{
            "front": "What is ... ?",
            "back": "A clear, direct answer in 1-2 sentences.",
            "hint": "Optional one-line hint (may be null)"
          }}
        ],
        "questions": [
          {{
            "question": "Question text?",
            "options": ["A", "B", "C", "D"],
            "correct_answer": 0,
            "explanation": "Why A is correct"
          }}
        ]
      }}
    ]
  }}
]

STUDY MATERIAL:
{text}"""


def _parse_ai_response(text: str) -> list:
    """Parse the AI response JSON."""
    # Find JSON in the response
    start = text.find('[')
    end = text.rfind(']') + 1
    if start == -1 or end == 0:
        logger.error('No JSON array found in AI response')
        return []

    try:
        return json.loads(text[start:end])
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse AI JSON: {e}')
        return []


def _generate_placeholder(text: str, num_topics: int, qpt: int) -> list:
    """
    Build study material without any AI key.

    Splits the text into roughly equal chunks, then for each chunk produces a
    subtopic with flashcards (sentence → key-term recall) and multiple-choice
    questions derived from the same sentences. Deterministic, no API required.
    """
    sentences = _sentences(text)
    if not sentences:
        sentences = [text.strip() or 'Sample study content']

    n_topics = max(1, min(num_topics, 4))
    chunk_size = max(1, len(sentences) // n_topics)
    # Build the pool of all derivable flashcards once so we can use them as
    # distractors across chunks, and to pad when a chunk is small.
    all_cards = [c for c in (_flashcard_from_sentence(s) for s in sentences) if c]
    topics = []

    for i in range(n_topics):
        start = i * chunk_size
        end = start + chunk_size if i < n_topics - 1 else len(sentences)
        chunk = sentences[start:end] or sentences[start:] or sentences
        if not chunk:
            chunk = [f'Topic {i + 1}']

        title = _topic_title(chunk, i)
        subtopic_title = _subtopic_title(chunk, i)
        flashcards = _flashcards_from_sentences(chunk, limit=max(4, min(8, qpt // 2)))
        questions = _questions_from_sentences(chunk, limit=min(qpt, 6), distractor_pool=all_cards)

        topics.append({
            'title': title,
            'description': chunk[0][:180],
            'subtopics': [{
                'title': subtopic_title,
                'description': chunk[0][:180],
                'flashcards': flashcards,
                'questions': questions,
            }],
        })

    return topics


def _topic_title(chunk: list, idx: int) -> str:
    """Pick a short, readable title from a chunk of sentences."""
    # Try to find a capitalised noun phrase near the start.
    first = chunk[0].strip()
    words = first.split()
    # Prefer the first meaningful capitalised run (e.g. "Calculus Fundamentals")
    cap_run = []
    for w in words[:12]:
        stripped = w.strip(',.:;')
        if stripped and stripped[0].isupper() and stripped.lower() not in {'the', 'a', 'an', 'and', 'or', 'but'}:
            cap_run.append(stripped)
            if len(cap_run) >= 4:
                break
        elif cap_run:
            break
    if cap_run:
        return f'Section {idx + 1}: ' + ' '.join(cap_run)
    # Fallback: first six words
    return f'Section {idx + 1}: ' + ' '.join(words[:6]).rstrip(',.:;')


def _subtopic_title(chunk: list, idx: int) -> str:
    """Short, readable subtopic title derived from the first sentence."""
    first = chunk[0].strip().rstrip('.:;')
    words = first.split()
    return ' '.join(words[:10]) + ('…' if len(words) > 10 else '')


def _sentences(text: str) -> list:
    """Split text into non-trivial sentences."""
    import re
    raw = re.split(r'(?<=[.!?])\s+', (text or '').strip())
    return [s.strip() for s in raw if len(s.strip()) > 12]


def _flashcards_from_sentences(sents: list, limit: int) -> list:
    """Turn sentences into front/back flashcards by extracting a keyword."""
    cards = []
    for s in sents[:limit * 2]:
        card = _flashcard_from_sentence(s)
        if card:
            cards.append(card)
        if len(cards) >= limit:
            break
    # Always give the user something: fall back to a whole-sentence recall card.
    while len(cards) < min(limit, len(sents)):
        idx = len(cards)
        if idx >= len(sents):
            break
        cards.append({
            'front': f'Recall the key point from line {idx + 1}.',
            'back': sents[idx][:300],
            'hint': None,
        })
    return cards


_STOPWORDS = {
    'the', 'a', 'an', 'and', 'or', 'but', 'is', 'are', 'was', 'were', 'be', 'been',
    'of', 'in', 'on', 'to', 'for', 'with', 'that', 'this', 'these', 'those',
    'as', 'at', 'by', 'from', 'it', 'its', 'into', 'than', 'then', 'also',
    'can', 'may', 'will', 'would', 'should', 'could', 'has', 'have', 'had',
}


def _flashcard_from_sentence(sentence: str) -> dict | None:
    """
    Turn a single sentence into a flashcard by blanking out the most informative word.

    Scores each non-stopword token for how "content-bearing" it is (capitalisation,
    length, digits), then blanks out the highest scorer.
    """
    import re
    words = sentence.split()
    if len(words) < 5:
        return None

    candidates = []
    for i, w in enumerate(words):
        stripped = re.sub(r'[^A-Za-z0-9\-]', '', w)
        if not stripped or stripped.lower() in _STOPWORDS:
            continue
        score = 0
        if stripped[0].isupper() and i > 0 and not stripped.isupper():
            score += 4
        if len(stripped) >= 8:
            score += 3
        elif len(stripped) >= 6:
            score += 2
        if any(ch.isdigit() for ch in stripped):
            score += 3
        # Domain verbs that signal a definition — favour the subject instead.
        if stripped.lower() in {'means', 'defined', 'called', 'refers', 'represents'}:
            score -= 2
        if score >= 2:
            candidates.append((score, i, stripped))

    if not candidates:
        return None

    candidates.sort(key=lambda t: (-t[0], t[1]))
    _, idx, term = candidates[0]
    masked = list(words)
    masked[idx] = '______'
    front = ' '.join(masked).strip().rstrip('.').rstrip(',')
    return {
        'front': front + '?' if not front.endswith('?') else front,
        'back': term,
        'hint': f'Starts with "{term[0]}" · {len(term)} letters',
    }


def _questions_from_sentences(sents: list, limit: int, distractor_pool: list | None = None) -> list:
    """Cheap multiple-choice generator for when no AI is available."""
    import random
    chunk_pool = [c for c in (_flashcard_from_sentence(s) for s in sents) if c]
    # Global pool lets us pull distractors from outside this chunk when we run out.
    global_pool = distractor_pool or chunk_pool

    questions = []
    for i, card in enumerate(chunk_pool[:limit]):
        correct = card['back']
        distractors = [c['back'] for c in global_pool if c['back'].lower() != correct.lower()]
        # Stable per-answer shuffle so tests are deterministic.
        rng = random.Random(hash(correct) & 0xFFFFFFFF)
        rng.shuffle(distractors)
        options = [correct] + distractors[:3]
        while len(options) < 4:
            options.append(f'Option {chr(ord("A") + len(options))}')
        rng.shuffle(options)
        questions.append({
            'question': card['front'].replace('______', '_____'),
            'options': options,
            'correct_answer': options.index(correct),
            'explanation': f'The correct answer is "{correct}" — drawn directly from the source material.',
        })

    # If we still didn't reach the target, synthesise simple true/false-style
    # questions from remaining sentences rather than emitting "placeholder N".
    remaining_sents = [s for s in sents if not any(s.startswith(c['front'].split('______')[0][:20]) for c in chunk_pool[:limit])]
    while len(questions) < limit and remaining_sents:
        s = remaining_sents.pop(0)
        lead = s.split('.')[0][:120]
        questions.append({
            'question': f'Which statement is true based on the source material?',
            'options': [
                lead,
                lead.replace(' is ', ' is not ') if ' is ' in lead else f'Not: {lead}',
                lead.replace(' the ', ' no ') if ' the ' in lead else f'None of the above',
                'None of the above',
            ][:4],
            'correct_answer': 0,
            'explanation': f'This is stated directly in the material: "{lead}".',
        })

    return questions


def ensure_flashcards_on_subtopic(subtopic: dict, source_sentences: list) -> list:
    """
    Ensure a subtopic has flashcards. If the AI response omitted them, derive
    flashcards from the subtopic's questions (Q → correct answer) or from the
    raw source sentences as a last resort.
    """
    cards = subtopic.get('flashcards') or []
    if cards:
        return cards

    # Derive from questions first — they're already vetted by the AI.
    for q in subtopic.get('questions') or []:
        options = q.get('options') or []
        idx = q.get('correct_answer', 0)
        if 0 <= idx < len(options):
            cards.append({
                'front': (q.get('question') or '').strip() or 'What is the key idea here?',
                'back': options[idx],
                'hint': None,
            })
        if len(cards) >= 6:
            break

    if not cards:
        cards = _flashcards_from_sentences(source_sentences, limit=6)

    return cards
